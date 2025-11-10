import asyncio
import base64
import gc
import io
import pickle
import cv2 # Cho khử nhiễu (OpenCV)
import numpy as np
import requests
import time
import PyPDF2 # Cho Watermark PDF
from io import BytesIO
from typing import Union
import openai
from fastapi import UploadFile, HTTPException
from google.api_core.exceptions import PermissionDenied
from nltk import PunktSentenceTokenizer
import json
from googletrans import Translator
from skimage.metrics import structural_similarity as ssim

from PIL import Image, ImageDraw, ImageFont, ImageOps # Cho Watermark ảnh
from pdf2image import convert_from_bytes
from PyPDF2.errors import PdfReadError
from pyvi.ViTokenizer import tokenize
from oauth2client.service_account import ServiceAccountCredentials
from pydrive.auth import GoogleAuth
from pydrive.drive import GoogleDrive
from sklearn.feature_extraction.text import TfidfVectorizer, CountVectorizer
from sklearn.metrics.pairwise import cosine_similarity
from sqlalchemy.exc import SQLAlchemyError
from starlette.responses import JSONResponse

from reportlab.pdfgen import canvas # Cho Watermark PDF
from reportlab.lib.pagesizes import letter
from reportlab.lib import colors
from reportlab.lib.units import inch

import zipfile # Để đọc file .docx, .pptx, .xlsx, .odt
import xml.etree.ElementTree as ET # (Fallback nếu lxml lỗi)
try:
    from lxml import etree as lxml_etree
    XML_PARSER = 'lxml'
except ImportError:
    logging.warning("Thư viện 'lxml' không tìm thấy. Sử dụng 'xml.etree.ElementTree' (chậm hơn).")
    lxml_etree = ET
    XML_PARSER = 'et'
    
# Thư viện cho các định dạng Office
from pptx import Presentation
import openpyxl
from odf.opendocument import load as load_odt
from odf.text import P
from odf.meta import DocumentStatistic

import logging
import bcrypt

from transformers import RobertaModel, AutoTokenizer
from elasticsearch import Elasticsearch
from docx import Document as DocxDocument
import re
from elasticsearch.helpers import bulk, BulkIndexError, scan
import string
from datetime import datetime, date, timedelta
import magic # Để kiểm tra mime type chính xác hơn
import os
from google.cloud import vision

from sqlalchemy import and_, or_, delete
from app.core.config import settings
from app.db.engine_default import session_scope
from app.models.models import Document, UserDocument, User, Organization, Category, PrivateDoc, StarredDoc, \
    FileUpload, Department, PublicLink, WorkflowInstance

elasticsearch_url = settings.ELASTICSEARCH_ENPOINT
google_application_credentials = settings.GOOGLE_APPLICATION_CREDENTIALS
os.environ["GOOGLE_APPLICATION_CREDENTIALS"] = google_application_credentials
model_path = settings.MODEL_PATH
model = settings.MODEL
client = Elasticsearch([elasticsearch_url])
tokenizer = AutoTokenizer.from_pretrained(model)
model_embedding = RobertaModel.from_pretrained(model)
openai.api_key = settings.OPENAI_API_KEY_EMBEDDING

logging.basicConfig(level=logging.DEBUG, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

def ensure_meta_data_index_exists():
    """Ensure meta_data index exists in Elasticsearch"""
    try:
        if not client.indices.exists(index="meta_data"):
            client.indices.create(
                index="meta_data",
                body={
                    "mappings": {
                        "properties": {
                            "document_id": {"type": "integer"},
                            "category_id": {"type": "integer"},
                            "document_number": {"type": "text"},
                            "issuing_authority": {"type": "text"},
                            "date_of_issuance": {"type": "text"},
                            "signature": {"type": "text"},
                            "agency_address": {"type": "text"}
                        }
                    }
                }
            )
            logging.info("Created meta_data index in Elasticsearch")
    except Exception as e:
        logging.error(f"Error ensuring meta_data index exists: {e}")

def get_mime_type(file_content: bytes) -> str:
    """Sử dụng 'magic' để xác định MIME type một cách an toàn từ nội dung byte."""
    try:
        mime = magic.Magic(mime=True)
        return mime.from_buffer(file_content)
    except Exception as e:
        logging.warning(f"Không thể xác định MIME type: {e}")
        return "application/octet-stream"

async def denoise_image(file_content: bytes, filename: str):
    """
    UC-39/UC-87: Khử nhiễu cho file ảnh trước OCR (Triển khai thật) 
    Sử dụng OpenCV fastNlMeansDenoisingColored.
    """
    mime_type = get_mime_type(file_content)
    if not mime_type.startswith('image/'):
        return {"content": file_content, "denoiseInfo": {"denoised": False, "message": "Không phải file ảnh, bỏ qua khử nhiễu."}}
        
    try:
        nparr = np.frombuffer(file_content, np.uint8)
        img_cv = cv2.imdecode(nparr, cv2.IMREAD_COLOR)
        if img_cv is None:
             raise ValueError("Không thể decode ảnh bằng OpenCV")
             
        logging.info(f"Đang khử nhiễu cho ảnh: {filename}")
        
        # Áp dụng thuật toán khử nhiễu
        denoised_img = cv2.fastNlMeansDenoisingColored(img_cv, None, 10, 10, 7, 21)
        
        # Chuyển ảnh đã xử lý về lại bytes
        # (Sử dụng đuôi file gốc để encode)
        file_extension = os.path.splitext(filename)[1]
        is_success, buffer = cv2.imencode(file_extension, denoised_img)
        if not is_success:
            raise ValueError(f"Không thể encode ảnh sang định dạng {file_extension}")
            
        denoised_content = buffer.tobytes()
        
        return {
            "content": denoised_content,
            "denoiseInfo": {
                "denoised": True,
                "modelUsed": "cv2.fastNlMeansDenoisingColored",
                "message": "Khử nhiễu ảnh thành công.",
                "original_size": len(file_content),
                "denoised_size": len(denoised_content)
            }
        }
    except Exception as e:
        logging.error(f"Lỗi khi khử nhiễu ảnh {filename}: {e}")
        return {"content": file_content, "denoiseInfo": {"denoised": False, "message": f"Lỗi khử nhiễu: {e}"}}

async def perform_ocr(file_content: bytes, filename: str):
    """
    UC-87: Thực hiện OCR (Wrapper)
    Trả về nội dung OCR và tổng số trang.
    Tận dụng các hàm `convert_file_to_content` và `convert_images_to_text` đã có.
    """
    try:
        # 1. Chuyển file bytes (PDF, DOCX, TXT, IMG) thành text hoặc list ảnh
        images_or_text, total_pages = convert_file_to_content(BytesIO(file_content), filename)
        
        # 2. Chuyển ảnh thành text (sử dụng Google Vision API)
        text = convert_images_to_text(images_or_text)
        
        if not text:
            logging.warning(f"OCR cho file {filename} không trả về nội dung.")
            
        return {"ocrContent": text or "", "total_pages": total_pages}
        
    except Exception as e:
        logging.error(f"Lỗi nghiêm trọng trong chuỗi OCR cho file {filename}: {e}")
        # Trả về 1 nếu lỗi
        return {"ocrContent": "", "total_pages": 1}
    
# *** BỔ SUNG HÀM HELPER NÀY (Cho Highlight) ***
def _find_context_and_highlight(full_text: str, segment: str, radius_chars: int = 150) -> str:
    """
    UC-39/UC-88: Tìm đoạn 'segment' trong 'full_text' và trả về context xung quanh
    với thẻ <b> highlight.
    """
    try:
        # Tìm vị trí (không phân biệt hoa thường, ký tự đặc biệt)
        # 1. Chuẩn hóa text
        clean_text = re.sub(r'\s+', ' ', full_text).lower()
        clean_segment = re.sub(r'\s+', ' ', segment).lower()
        
        index = clean_text.find(clean_segment)
        if index == -1:
            return segment # Trả về segment gốc nếu không tìm thấy context
            
        # 2. Lấy context dựa trên index đã tìm thấy
        start = max(0, index - radius_chars)
        end = min(len(full_text), index + len(segment) + radius_chars)
        
        # 3. Lấy chuỗi gốc (có dấu, hoa thường)
        context_text = full_text[start:end]
        
        # 4. Tìm lại segment gốc trong context để highlight
        # (Tìm cách an toàn, không phân biệt hoa thường)
        segment_regex = re.compile(re.escape(segment), re.IGNORECASE)
        highlighted_context = segment_regex.sub(lambda m: f"<b>{m.group(0)}</b>", context_text, count=1)
        
        # Thêm "..." nếu bị cắt
        if start > 0:
            highlighted_context = "... " + highlighted_context
        if end < len(full_text):
            highlighted_context = highlighted_context + " ..."
            
        return highlighted_context
        
    except Exception as e:
        logging.error(f"Lỗi khi highlight (UC-88): {e}")
        return f"<b>{segment}</b>" # Fallback

async def check_duplicates_by_content(
    ocr_content: str, 
    organization_id: int, 
    user_id: int, 
    similarity_threshold: float = 0.30
):
    """
    UC-39/UC-88: Kiểm tra trùng lặp dựa trên tương đồng ngữ nghĩa (Triển khai thật 100%)
    """
    if not ocr_content.strip():
        logging.warning("Nội dung OCR rỗng, bỏ qua kiểm tra trùng lặp.")
        return {"isDuplicate": False}
        
    try:
        # 1. Xác định model và index
        is_openai_org = check_orgnization_open_ai(user_id)
        index_name = "search_openai" if is_openai_org else "demo_cau"
        
        # 2. Vector hóa nội dung mới
        if is_openai_org:
            headers = { "Authorization": f"Bearer {settings.OPENAI_API_KEY_EMBEDDING}" }
            query_vector = get_openai_embeddings_to_search([ocr_content], headers)[0]
        else:
            query_vector = embed_str([tokenize(ocr_content)])[0]
            
        # 3. Xây dựng truy vấn Elasticsearch
        score_threshold = (similarity_threshold + 1.0)
        
        script_query = {
            "bool": {
                "must": {
                    "script_score": {
                        "query": { "match_all": {} },
                        "script": {
                            "source": "cosineSimilarity(params.query_vector, 'title_vector') + 1.0",
                            "params": {"query_vector": query_vector}
                        },
                        "min_score": score_threshold
                    }
                },
                "filter": [
                    {"exists": {"field": "title_vector"}},
                    # --- TRIỂN KHAI YÊU CẦU 1 ---
                    # (Filter theo organization_id)
                    {"term": {"organization_id": organization_id}}
                ]
            }
        }
        
        # 4. Thực hiện tìm kiếm
        response = client.search(
            index=index_name,
            body={
                "query": script_query,
                "_source": ["document_id", "title"], # Lấy content chunk (đang bị lưu vào trường 'title')
                "size": 1, 
                "sort": [{"_score": {"order": "desc"}}]
            },
            ignore=[400, 404]
        )
        
        hits = response.get("hits", {}).get("hits", [])
        
        if not hits:
            return {"isDuplicate": False}
            
        # 5. Xử lý kết quả
        top_hit = hits[0]
        similarity_score = top_hit["_score"] - 1.0 
        duplicate_doc_id = top_hit["_source"]["document_id"]
        
        # Lấy content của chunk bị trùng (đang lưu trong trường 'title' của ES)
        duplicate_sentence_chunk = top_hit["_source"]["title"] 

        with session_scope() as session:
            duplicate_doc = session.query(Document).filter(Document.id == duplicate_doc_id).first()
            if not duplicate_doc:
                 return {"isDuplicate": False} 

            # --- TRIỂN KHAI YÊU CẦU 2 ---
            # (Highlight đoạn trùng) 
            # Tìm chunk trùng lặp (`duplicate_sentence_chunk`)
            # bên trong toàn bộ nội dung OCR mới (`ocr_content`)
            highlighted_segment = _find_context_and_highlight(
                full_text=ocr_content, 
                segment=duplicate_sentence_chunk
            )

            return {
                "isDuplicate": True,
                "duplicateData": {
                    "similarity": f"{similarity_score * 100:.1f}%",
                    "highlight": highlighted_segment, # Trả về đoạn highlight
                    "existingDocument": {
                        "name": duplicate_doc.title,
                        "id": duplicate_doc.id,
                        "owner": duplicate_doc.created_by
                    }
                }
            }
        
    except Exception as e:
        logging.error(f"Lỗi khi kiểm tra trùng lặp nội dung (UC-88): {e}")
        return {"isDuplicate": False, "message": str(e)}

async def suggest_metadata_from_content(ocr_content: str, filename: str):
    """
    UC-73: Gợi ý metadata từ nội dung (Triển khai thật)
    Sử dụng OpenAI (hàm `call_openai` đã có) để trích xuất KVP, title, tags.
    """
    if not ocr_content.strip():
        logging.warning("Nội dung OCR rỗng, sử dụng tên file làm tiêu đề.")
        return {
            "suggestedMetadata": { "title": filename.split('.')[0], "tags": [], "key_values": {} },
            "warnings": []
        }
        
    try:
        # Xây dựng prompt chi tiết theo UC-73
        prompt = (
            "Phân tích nội dung văn bản sau đây và trả về một đối tượng JSON DUY NHẤT. "
            "KHÔNG GIẢI THÍCH, CHỈ TRẢ VỀ JSON.\n"
            "Nội dung văn bản:\n"
            f"---BEGIN CONTENT---\n{ocr_content[:2000]}\n---END CONTENT---\n\n" # Lấy 2000 ký tự đầu
            "Đối tượng JSON phải có cấu trúc sau:\n"
            "{\n"
            "  \"title\": \"(Gợi ý một tiêu đề ngắn gọn, súc tích cho văn bản)\",\n"
            "  \"summary\": \"(Gợi ý một đoạn tóm tắt 2-3 câu về nội dung chính)\",\n"
            "  \"tags\": [\"(từ khóa 1)\", \"(từ khóa 2)\", \"(từ khóa 3)\"],\n"
            "  \"key_values\": {\n"
            "    \"so_hieu\": \"(Trích xuất Số hiệu/Số văn bản, ví dụ: 123/QĐ-BGDĐT)\",\n"
            "    \"ngay_ban_hanh\": \"(Trích xuất Ngày ban hành, định dạng dd/mm/yyyy)\",\n"
            "    \"co_quan_ban_hanh\": \"(Trích xuất Cơ quan ban hành)\",\n"
            "    \"nguoi_ky\": \"(Trích xuất Người ký, nếu có)\",\n"
            "    \"so_luong\": (Trích xuất số lượng dạng SỐ, nếu có, ví dụ: 500),\n"
            "    \"ngay_hieu_luc\": \"(Trích xuất Ngày hiệu lực, nếu có, dd/mm/yyyy)\"\n"
            "  }\n"
            "}\n"
            "Nếu không tìm thấy thông tin cho trường nào, hãy trả về null cho trường đó."
        )
        
        # (Giả định hàm `call_openai` đã có và trả về CHUỖI JSON)
        response_text = call_openai(prompt)
        
        if not response_text:
            raise ValueError("OpenAI không trả về kết quả.")
            
        # Phân tích chuỗi JSON
        generated_data = json.loads(response_text)
        
        # Xử lý giá trị thiếu (UC-73)
        key_values = generated_data.get("key_values", {})
        warnings = []
        
        # Điền null hoặc warning cho các trường bị thiếu
        for key, value in key_values.items():
            if value is None:
                warnings.append({
                    "field": key,
                    "message": f"Không tìm thấy thông tin cho '{key}'."
                })
        
        # Đảm bảo các trường chính luôn tồn tại
        suggested_title = generated_data.get("title") or filename.split('.')[0]
        tags = generated_data.get("tags", [])
        
        return {
            "suggestedMetadata": {
                "title": suggested_title,
                "tags": tags,
                "summary": generated_data.get("summary", ""),
                "key_values": key_values,
                "category": 1 # Mặc định "Hành chính"
            },
            "warnings": warnings
        }
        
    except Exception as e:
        logging.error(f"Lỗi khi gợi ý metadata (UC-73): {e}")
        return {
            "suggestedMetadata": { "title": filename.split('.')[0], "tags": [], "key_values": {} },
            "warnings": [{"field": "System", "message": f"Lỗi trích xuất AI: {e}"}]
        }

def check_data_conflicts(key_values: dict):
    """
    UC-39/UC-73: Kiểm tra mâu thuẫn dữ liệu (Triển khai thật)
    Kiểm tra số lượng âm và ngày vượt quá 20/08/2025.
    """
    conflicts = []
    
    # Rule 1: Kiểm tra "số lượng" âm (UC-39, UC-73)
    # (Lưu ý: Tên trường 'so_luong' phải khớp với kết quả từ suggest_metadata)
    so_luong = key_values.get("so_luong")
    if so_luong is not None:
        try:
            # Thử ép kiểu về số, phòng trường hợp AI trả về "500"
            numeric_so_luong = float(so_luong)
            if numeric_so_luong < 0:
                conflicts.append({
                    "field": "so_luong",
                    "value": so_luong,
                    "message": f"Giá trị mâu thuẫn: Số lượng ({so_luong}) không thể là số âm."
                })
        except (ValueError, TypeError):
            pass # Bỏ qua nếu không phải là số

    # Rule 2: Kiểm tra "ngày ban hành" vượt quá 20/08/2025 (UC-39, UC-73)
    ngay_ban_hanh_str = key_values.get("ngay_ban_hanh")
    if ngay_ban_hanh_str:
        try:
            # (AI có thể trả về nhiều định dạng, ở đây xử lý dd/mm/yyyy)
            doc_date = datetime.strptime(ngay_ban_hanh_str, "%d/%m/%Y")
            # (Tham chiếu Thông tư 05/2025)
            cutoff_date = datetime(2025, 8, 20) 
            
            if doc_date > cutoff_date:
                conflicts.append({
                    "field": "ngay_ban_hanh",
                    "value": ngay_ban_hanh_str,
                    "message": f"Giá trị mâu thuẫn: Ngày ban hành ({ngay_ban_hanh_str}) vượt quá mốc thời gian (20/08/2025)."
                })
        except (ValueError, TypeError):
            # Bỏ qua nếu định dạng ngày không hợp lệ
            logging.warning(f"Không thể parse ngày '{ngay_ban_hanh_str}' để kiểm tra mâu thuẫn.")
            
    return {"conflicts": conflicts}

async def embed_watermark_preview(file_content: bytes, filename: str):
    """
    UC-39: Chỉ là bước xử lý giả lập, trả về thông báo cho Giai đoạn 1.
    Watermark thật sẽ được nhúng ở Giai đoạn 2 (embed_watermark_final).
    """
    await asyncio.sleep(0.1) # Giả lập
    return {
        "watermarkInfo": {
            "success": True,
            "message": "Sẵn sàng nhúng watermark 'Confidential - INNOTECH' (nếu cần).",
        }
    }

# *** THAY THẾ HÀM NÀY ***
def upload_file_to_drive(file_content: bytes, unique_filename: str):  
    """
    Tải file (dưới dạng bytes) lên Google Drive và lưu tạm 1 bản local.
    """
    file_directory = settings.UPLOAD_FILE_DIRECTORY
    os.makedirs(file_directory, exist_ok=True)
    
    # Lưu file local (để các tiến trình background khác có thể truy cập nếu cần)
    local_file_path = os.path.join(file_directory, unique_filename)
    
    try:
        with open(local_file_path, 'wb') as f:
            f.write(file_content)
    except Exception as e:
        logging.error(f"Không thể ghi file local {local_file_path}: {e}")
        raise

    gauth = GoogleAuth()
    scope = ['https://www.googleapis.com/auth/drive']
    gauth.credentials = ServiceAccountCredentials.from_json_keyfile_name(
        google_application_credentials, scope)
    drive = GoogleDrive(gauth)

    uploaded_file = drive.CreateFile({'title': unique_filename})
    uploaded_file.SetContentFile(local_file_path)
    uploaded_file.Upload()

    if uploaded_file.get('id'):
        uploaded_file.InsertPermission({
            'type': 'anyone',
            'value': 'anyone',
            'role': 'reader'
        })
        file_id = uploaded_file.get('id')
        file_url = uploaded_file['alternateLink']
        
        # Trả về đường dẫn local để hàm insert_file_upload (nếu cần)
        return file_id, file_url, local_file_path
    else:
        # Xóa file local nếu upload lỗi
        if os.path.exists(local_file_path):
            os.remove(local_file_path)
        raise Exception("File upload failed to Google Drive")

def format_document_number(doc_number):
    if not doc_number:
        return None

    # Regex to match the format, including the new case
    match = re.match(r'^(\d{1,3})(\d{4})([A-Z\-Đ]+)$', doc_number)
    if match:
        serial_number = match.group(1)
        year = match.group(2)
        authority_abbr = match.group(3)

        # Ensure that abbreviation is formatted as TT-ABC or NĐ-CP
        if len(authority_abbr) > 2:
            formatted_abbr = f"{authority_abbr[:2]}-{authority_abbr[2:]}"
        else:
            formatted_abbr = authority_abbr

        return f"{serial_number}/{year}/{formatted_abbr}"
    else:
        # If format does not match, return the original
        return doc_number


def call_openai(prompt):
    try:
        # Truncate the prompt to 1000 characters if it exceeds this length
        truncated_prompt = prompt[:1000]
        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": "You are a helpful assistant."},
                {"role": "user", "content": truncated_prompt}
            ]
        )
        response_text = response.choices[0].message['content'].strip()

        # Log or print the full response for debugging
        print(f"OpenAI response: {response_text}")

        # Use regex to extract JSON content
        json_pattern = re.compile(r'\{.*\}', re.DOTALL)
        match = json_pattern.search(response_text)
        if not match:
            return json.dumps({
                'document_number': "",
                'issuing_authority': "",
                'date_of_issuance': "",
                'signature': "",
                'agency_address': ""
        })  # Re
        response_json = match.group(0).strip()
        data = json.loads(response_json)

        # Assign a default value if document_number is missing or empty
        if not data.get('document_number'):
            data['document_number'] = ""

        # Format the document_number if present
        if 'document_number' in data:
            data['document_number'] = format_document_number(data['document_number'])

        if not data.get('issuing_authority'):
            data['issuing_authority'] = ""
        # Convert issuing_authority to uppercase if present
        if data.get('issuing_authority'):
            data['issuing_authority'] = data.get('issuing_authority').upper()

        if not data.get('date_of_issuance'):
            data['date_of_issuance'] = ""

        if not data.get('signature'):
            data['signature'] = ""

        if not data.get('agency_address'):
            data['agency_address'] = ""

        # Ensure all fields are included and return them
        return json.dumps({
            'document_number': data.get('document_number'),
            'issuing_authority': data.get('issuing_authority'),
            'date_of_issuance': data.get('date_of_issuance'),
            'signature': data.get('signature'),
            'agency_address': data.get('agency_address')
        })  # Return the JSON as a string
        

    except Exception as e:
        print(f"OpenAI API call error: {e}")
        return None


def get_spliter():
    with open(model_path, 'rb') as fs:
        punkt_param = pickle.load(fs)

    punkt_param.sent_starters = {}
    abbrev_types = ['g.m.t', 'e.g', 'dr', 'dr', 'vs', "000", 'mr', 'mrs', 'prof', 'inc', 'tp', 'ts', 'ths',
                    'th', 'vs', 'tp', 'k.l', 'a.w.a.k.e', 't', 'a.i', '</i', 'g.w',
                    'ass',
                    'u.n.c.l.e', 't.e.s.t', 'ths', 'd.c', 've…', 'ts', 'f.t', 'b.b', 'z.e', 's.g', 'm.p',
                    'g.u.y',
                    'l.c', 'g.i', 'j.f', 'r.r', 'v.i', 'm.h', 'a.s', 'bs', 'c.k', 'aug', 't.d.q', 'b…', 'ph',
                    'j.k', 'e.l', 'o.t', 's.a']
    abbrev_types.extend(string.ascii_uppercase)
    for abbrev_type in abbrev_types:
        punkt_param.abbrev_types.add(abbrev_type)
    for abbrev_type in string.ascii_lowercase:
        punkt_param.abbrev_types.add(abbrev_type)
    return PunktSentenceTokenizer(punkt_param)


sent_tokenizer = get_spliter()


def embed_str(str):
    # Initialize the tokenizer
    encoded_input = tokenizer(str, padding=True, truncation=True, return_tensors='pt')

    # Pass the encoded input through the model
    output = model_embedding(**encoded_input)

    # Extract the embeddings from the model's output
    embeddings = output.last_hidden_state.mean(dim=1).detach().numpy()

    return embeddings


def delete_file_in_drive(file_id):
    gauth = GoogleAuth()
    # Load the Service Account credentials
    scope = ['https://www.googleapis.com/auth/drive']
    gauth.credentials = ServiceAccountCredentials.from_json_keyfile_name(
        google_application_credentials, scope)
    drive = GoogleDrive(gauth)

    # Create a file instance and delete it
    file_drive = drive.CreateFile({'id': file_id})
    file_drive.Delete()


def validate_file(file: UploadFile):
    file_size = os.fstat(file.file.fileno()).st_size
    max_file_size = 1000 * 1024 * 1024
    if file_size > max_file_size:
        raise HTTPException(status_code=400, detail="The file size exceeds the limit!")

    return True


def convert_file_to_content(file: BytesIO, filename: str) -> (Union[str, list], int):
    """
    Chuyển đổi file (bytes) sang nội dung (text hoặc list ảnh)
    VÀ trả về tổng số trang (total_pages) thực tế.
    """
    print(f"Bắt đầu convert_file_to_content cho: {filename}")
    file_format = filename.split(".")[-1].lower()
    total_pages = 1 # Mặc định là 1 (cho các file không thể đếm)
    text_content = "" # Nội dung text
    images_content = None # Nội dung ảnh (cho PDF)
    
    file.seek(0) # Đảm bảo con trỏ file ở đầu

    if file_format == "pdf":
        pdf_bytes = file.read()
        file.seek(0)
        
        # 1. Dùng pdf2image để lấy ảnh (cho OCR)
        images_content = convert_from_bytes(pdf_bytes)
        
        # 2. Dùng PyPDF2 để đếm trang (chính xác)
        try:
            pdf_reader = PyPDF2.PdfReader(file)
            total_pages = len(pdf_reader.pages)
        except Exception as e_count:
            logging.warning(f"Không thể đếm trang PDF: {e_count}. Dùng count của pdf2image.")
            total_pages = len(images_content)
            
        # (images_content sẽ được xử lý bởi `convert_images_to_text`)
        return images_content, total_pages

    elif file_format in ["docx", "doc"]:
        # Lấy text
        docx_document = DocxDocument(file)
        text_content = "\n".join([paragraph.text for paragraph in docx_document.paragraphs])
        
        # Đếm trang (Cách 1: Đọc metadata)
        # DOCX là file ZIP. Chúng ta sẽ đọc file 'docProps/app.xml'
        try:
            file.seek(0)
            with zipfile.ZipFile(file, 'r') as docx_zip:
                app_xml = docx_zip.read('docProps/app.xml')
                # (Sử dụng lxml nếu có)
                root = lxml_etree.fromstring(app_xml)
                # Namespace cho docProps
                ns = {'ep': 'http://schemas.openxmlformats.org/officeDocument/2006/extended-properties'}
                
                # Tìm thẻ <Pages>
                pages_tag = root.find('.//ep:Pages', namespaces=ns)
                if pages_tag is not None and pages_tag.text.isdigit():
                    total_pages = int(pages_tag.text)
                else:
                    # (Cách 2: Fallback - Đếm ngắt trang thủ công - phức tạp, tạm thời giữ 1)
                    logging.warning(f"Không tìm thấy <Pages> tag trong app.xml cho file {filename}. Tạm tính 1 trang.")
                    total_pages = 1 # (Việc đếm thủ công (manual break) rất không chính xác)
                    
        except Exception as e_docx:
            logging.error(f"Lỗi khi đếm trang DOCX {filename}: {e_docx}. Tạm tính 1 trang.")
            total_pages = 1
            
        return text_content, total_pages

    elif file_format == "txt":
        try:
            text = file.read().decode('utf-8')
        except UnicodeDecodeError:
            try:
                file.seek(0)
                text = file.read().decode('iso-8859-1')
            except UnicodeDecodeError:
                raise HTTPException(status_code=400, detail="Mã hóa tệp không được hỗ trợ.")
        total_pages = 1 
        return text, total_pages

    elif file_format == "xlsx":
        try:
            import pandas as pd
            xlsx = pd.ExcelFile(file)
            text_content = ""
            for sheet_name in xlsx.sheet_names:
                df = xlsx.parse(sheet_name, header=None)
                df = df.astype(str)
                text_content += df.apply(lambda x: ' '.join(x), axis=1).str.cat(sep='\n')
                text_content += "\n"
        except ImportError:
            raise HTTPException(status_code=500, detail="Thiếu 'pandas' hoặc 'openpyxl'.")
            
        # Đếm trang (Excel không có "trang", ta đếm "worksheets")
        try:
            file.seek(0)
            wb = openpyxl.load_workbook(file, read_only=True)
            total_pages = len(wb.sheetnames) # Đếm số lượng worksheets
        except Exception as e_xlsx:
            logging.error(f"Lỗi khi đếm sheet XLSX {filename}: {e_xlsx}. Tạm tính 1 trang.")
            total_pages = 1
            
        return text_content, total_pages

    elif file_format == "pptx":
        from pptx import Presentation
        file.seek(0)
        ppt = Presentation(file)
        text_content = ""
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content += shape.text + "\n"
                    
        # Đếm trang (PPTX đếm "slides")
        total_pages = len(ppt.slides) # Đếm số lượng slides
        
        return text_content, total_pages

    elif file_format in ["odt", "rtf"]:
        from pyth.plugins.rtf15.reader import Rtf15Reader
        from pyth.plugins.plaintext.writer import PlaintextWriter
        doc = Rtf15Reader.read(file)
        text = PlaintextWriter.write(doc).getvalue()
        total_pages = 1
        return text, total_pages

    else:
        raise HTTPException(status_code=400, detail="Định dạng tệp không hợp lệ!")


def convert_images_to_text(images):
    print("Starting convert_images_to_text")
    text = ""
    client = vision.ImageAnnotatorClient()

    try:
        if isinstance(images, list):  # Hình ảnh từ tệp PDF
            for image in images:
                # Convert image to bytes
                byte_arr = io.BytesIO()
                image.save(byte_arr, format='PNG')
                encoded_image = base64.b64encode(byte_arr.getvalue()).decode()

                image = vision.Image(content=encoded_image)
                response = client.text_detection(image=image)
                texts = response.text_annotations
                if texts:
                    text += texts[0].description
        elif isinstance(images, str):  # Văn bản từ tệp DOCX hoặc TXT
            text = images

        # Loại bỏ ký tự đặc biệt, trừ dấu chấm và dấu phẩy
        special_chars = string.punctuation.replace('.', '').replace(',', '')
        cleaned_text = text.replace("\n", " ").translate(str.maketrans('', '', special_chars))

        return cleaned_text

    except PermissionDenied as e:
        logging.error(f"PermissionDenied: {e.message}")
        raise HTTPException(status_code=403,
                            detail="Cloud Vision API is not enabled or permission is denied. Please enable the API and try again.")

    except Exception as e:
        logging.error(f"Error in convert_images_to_text: {e}")
        raise HTTPException(status_code=500, detail="An error occurred while processing the images.")


def split_chunks_into_sentences(chunks):
    if isinstance(chunks, str):
        chunks = [chunks]  # treat the entire text as a single chunk

    sentences = []
    for chunk in chunks:
        sentences.extend(re.split('(?<=[.!?]) +', chunk))
    return sentences


def check_duplicate_document(title: str, category_id: int) -> str:
    with session_scope() as session:
        suffix = 1
        base_name, extension = os.path.splitext(title)
        new_title = title
        while True:
            existing_document = session.query(Document).filter(Document.title == new_title,
                                                               Document.category_id == category_id).first()
            if existing_document is None:
                break
            new_title = f"{base_name}({suffix}){extension}"
            suffix += 1
    return new_title


def split_text_into_chunks(text):
    chunks = re.split(r'\n\n|\.\s*\n', text)
    return chunks

# *** THAY THẾ HÀM NÀY (Thêm organization_id) ***
def create_embeddings_and_store(texts: str, filename: str, category_id: int, document_id: int, type_doc: str, organization_id: int):
    print("Starting create_embeddings_and_store")
    # Xử lý trường hợp texts là None hoặc không phải là iterable
    if texts is None:
        logging.warning(f"texts is None for document {document_id}, skipping embeddings creation")
        return False
    
    text = ""
    if isinstance(texts, str):
        text = texts
    elif isinstance(texts, (list, tuple)):
        for i in texts:
            if i:
                text += str(i)
    else:
        try:
            for i in texts:
                text += str(i)
        except TypeError:
            logging.warning(f"texts is not iterable for document {document_id}, skipping embeddings creation")
            return False

    sentences = text.replace("\n", " ")
    data = sent_tokenizer.sentences_from_text(sentences)

    try:
        # Truyền organization_id vào
        result = preprocessing_indexing_elasticsearch(data, filename, category_id, document_id, type_doc, organization_id)
        index_batch(result) # (Hàm index_batch không cần đổi, vì nó chỉ đẩy 'result' đi)
        print("Indexing completed successfully")
    except Exception as e:
        print(f"An error occurred during indexing: {e}")
        return False

    del data
    del result
    gc.collect()
    print("Finished create_embeddings_and_store")
    return True


def text_to_wordlist(text):
    arr_input = text
    arr_input = re.sub(r'[?|$|.,…|!|:;\\*\-+=^/\'\"\”@`~#%&\]\[\{\}–—­­­­­_\()]', r' ', arr_input)
    arr_input = arr_input.replace('tnghị quyết', 'nghị quyết')
    arr_input = arr_input.replace('tướngphó', 'tướng phó')
    arr_input = arr_input.replace('\n', ' ')
    arr_input = " ".join(arr_input.split())
    words = arr_input.lower().split()
    return words


def embed_text(batch_text):
    encoded_input = tokenizer(batch_text, padding=True, truncation=True, return_tensors='pt')
    output = model_embedding(**encoded_input)
    embeddings = output.last_hidden_state.mean(dim=1).detach().numpy()
    return embeddings


def index_batch(docs):
    batch_size = 5
    print("start_index_batch")
    for i in range(0, len(docs), batch_size):
        batch_docs = docs[i:i + batch_size]
        requests = []
        titles = [tokenize(doc["title"]) for doc in batch_docs]
        title_vectors = embed_text(titles)
        for j, doc in enumerate(batch_docs):
            request = doc
            request["_op_type"] = "index"
            request["_index"] = "demo_cau"
            request["title_vector"] = title_vectors[j]
            requests.append(request)
        try:
            print("start_bulk")
            bulk(client, requests)
        except BulkIndexError as e:
            print(f"Bulk index error: {e}")

        del requests
        del titles
        del title_vectors
        gc.collect()


# *** THAY THẾ HÀM NÀY (Thêm organization_id) ***
def preprocessing_indexing_elasticsearch(arr_input, filename, category_id, document_id, type_doc, organization_id: int):
    result = []
    dem = 0
    for i in range(len(arr_input)):
        arr_input[i] = (text_to_wordlist(arr_input[i].lower()))
        arr_input[i] = " ".join(arr_input[i])
        if len(arr_input[i]) < 12:
            continue
        dem = dem + 1
        result.append({
            'id': dem, 
            'title': arr_input[i], # 'title' thực ra là 'content chunk'
            'doc_id': filename, 
            'category_id': category_id,
            'document_id': document_id, 
            'type': type_doc,
            'organization_id': organization_id # <-- TRƯỜNG MỚI QUAN TRỌNG
        })
    return result

async def embed_watermark_final(file_content: bytes, filename: str, confidentiality: str):
    """
    UC-39 (Step 3.2): Nhúng watermark THẬT sự 
    - PUBLIC: Nhúng logo INNOTECH (dạng ảnh)
    - INTERNAL/LOCKED: Nhúng text "Confidential - INNOTECH"
    """
    mime_type = get_mime_type(file_content)
    
    # 1. XỬ LÝ FILE ẢNH (JPG, PNG)
    if mime_type.startswith('image/'):
        try:
            image = Image.open(BytesIO(file_content)).convert("RGBA")
            txt_overlay = Image.new("RGBA", image.size, (255, 255, 255, 0))
            draw = ImageDraw.Draw(txt_overlay)
            
            # Giả sử font đã có trong thư mục (cần cho Python 3.8-slim)
            # Bạn cần đảm bảo file "arial.ttf" tồn tại trong thư mục Be-python/app
            font_path = os.path.join(os.path.dirname(__file__), '..', 'arial.ttf')
            if not os.path.exists(font_path):
                # Fallback nếu không có font
                font_path = "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf" # Đường dẫn phổ biến trên Linux
                if not os.path.exists(font_path):
                    logging.warning("Không tìm thấy font. Watermark text có thể không hoạt động.")
                    font = ImageFont.load_default()
                else:
                    font = ImageFont.truetype(font_path, size=40)
            else:
                 font = ImageFont.truetype(font_path, size=40)

            if confidentiality == "PUBLIC":
                # UC-39: Nhúng logo INNOTECH
                # (Giả sử bạn có file logo-teal.png trong thư mục Be-python/)
                logo_path = os.path.join(os.path.dirname(__file__), '..', 'logo-teal.png')
                if os.path.exists(logo_path):
                    logo = Image.open(logo_path).convert("RGBA")
                    # Resize logo (ví dụ: bằng 1/5 chiều rộng ảnh)
                    logo_width = image.width // 5
                    logo_height = int(logo.height * (logo_width / logo.width))
                    logo = logo.resize((logo_width, logo_height), Image.Resampling.LANCZOS)
                    
                    # Đặt opacity cho logo
                    logo = Image.eval(logo, lambda p: min(p, 100)) # 100/255 opacity (khoảng 40%)
                    
                    # Vị trí (ví dụ: góc dưới bên phải)
                    position = (image.width - logo_width - 20, image.height - logo_height - 20)
                    txt_overlay.paste(logo, position, logo)
                else:
                    logging.warning(f"Không tìm thấy file logo tại {logo_path}")

            elif confidentiality == "INTERNAL" or confidentiality == "LOCKED":
                # UC-39: Nhúng text ẩn "Confidential - INNOTECH"
                text = "Confidential - INNOTECH"
                fill_color = (255, 0, 0, 60) # Màu đỏ, mờ (alpha 60/255)
                
                # Tính toán vị trí (xoay và trung tâm)
                draw.text((0,0), text, font=font, fill=fill_color)
                text_bbox = draw.textbbox((0, 0), text, font=font)
                text_width = text_bbox[2] - text_bbox[0]
                text_height = text_bbox[3] - text_bbox[1]
                
                # Tạo ảnh text đã xoay
                text_image = Image.new('RGBA', (text_width, text_height), (255, 255, 255, 0))
                text_draw = ImageDraw.Draw(text_image)
                text_draw.text((0, 0), text, font=font, fill=fill_color)
                rotated_text = text_image.rotate(45, expand=True, fillcolor=(0,0,0,0))
                
                # Dán ảnh đã xoay vào trung tâm overlay
                x = (image.width - rotated_text.width) // 2
                y = (image.height - rotated_text.height) // 2
                txt_overlay.paste(rotated_text, (x, y), rotated_text)
            
            # Kết hợp ảnh gốc và overlay
            out_image = Image.alpha_composite(image, txt_overlay)
            
            byte_arr = BytesIO()
            # Lưu ảnh về lại định dạng gốc (nếu là PNG) hoặc JPG (nếu là TIFF/...)
            file_format = "PNG" if mime_type == "image/png" else "JPEG"
            out_image.convert("RGB").save(byte_arr, format=file_format)
            return byte_arr.getvalue()

        # 2. XỬ LÝ FILE PDF
        except Exception as e_img:
            logging.error(f"Lỗi khi nhúng watermark ảnh {filename}: {e_img}")
            return file_content # Trả về file gốc nếu lỗi
            
    elif mime_type == 'application/pdf':
        try:
            watermark_bytes = BytesIO()
            c = canvas.Canvas(watermark_bytes, pagesize=letter)
            
            if confidentiality == "PUBLIC":
                # UC-39: Nhúng logo INNOTECH
                logo_path = os.path.join(os.path.dirname(__file__), '..', 'logo-teal.png')
                if os.path.exists(logo_path):
                    # reportlab có thể nhận đường dẫn file trực tiếp
                    # Vị trí góc dưới bên phải
                    c.drawImage(logo_path, letter[0] - 1.2*inch, 0.5*inch, width=1*inch, preserveAspectRatio=True, mask='auto', anchor='c')
                else:
                    logging.warning(f"Không tìm thấy file logo tại {logo_path}")
                    
            elif confidentiality == "INTERNAL" or confidentiality == "LOCKED":
                # UC-39: Nhúng text ẩn "Confidential - INNOTECH"
                text = "Confidential - INNOTECH"
                c.setFont("Helvetica", 50)
                c.setFillColor(colors.red, alpha=0.15) # Màu đỏ, rất mờ
                c.saveState()
                c.translate(letter[0]/2, letter[1]/2) # Trung tâm
                c.rotate(45) # Xoay 45 độ
                c.drawCentredString(0, 0, text)
                c.restoreState()
            
            c.save()
            watermark_bytes.seek(0)
            
            # Đọc watermark vừa tạo
            watermark_pdf = PyPDF2.PdfReader(watermark_bytes)
            if not watermark_pdf.pages:
                raise ValueError("Tạo trang watermark thất bại.")
            watermark_page = watermark_pdf.pages[0]
            
            output_writer = PyPDF2.PdfWriter()
            input_pdf = PyPDF2.PdfReader(BytesIO(file_content))
            
            # Lặp qua từng trang của PDF gốc và merge watermark
            for page in input_pdf.pages:
                page.merge_page(watermark_page)
                output_writer.add_page(page)
                
            output_bytes = BytesIO()
            output_writer.write(output_bytes)
            return output_bytes.getvalue()

        except (PdfReadError, Exception) as e_pdf:
            logging.error(f"Lỗi khi nhúng watermark PDF {filename}: {e_pdf}. Trả về file gốc.")
            return file_content

    # 3. XỬ LÝ DOCX (Vẫn log cảnh báo)
    elif "wordprocessingml" in mime_type:
        logging.warning(f"Watermark cho DOCX ({filename}) không được hỗ trợ tự động trên server. Bỏ qua...")
        return file_content

    else:
        # Các định dạng khác (MP4, MP3...)
        return file_content

async def insert_document(filename: str, user: User, category_id: int, description: str,
                          file_type: str, file_id: str, file_path: str, storage_capacity: int,
                          storage_unit: str, photo_id: str, total_pages: int,
                          # --- Tham số mới ---
                          final_metadata: dict,
                          status_str: str, # "DRAFT"
                          version: str): # "1.0"
    with session_scope() as session:
        created_by = user.email if user else None
        
        # Trích xuất metadata từ object mới
        title = final_metadata.get('title', filename)
        access_type_str = final_metadata.get('access_type', 'private')
        confidentiality = final_metadata.get('confidentiality', 'INTERNAL')
        
        # --- ÁNH XẠ LOGIC (Rất quan trọng) ---
        
        # 1. Ánh xạ Status (String) sang Status (Integer)
        # Dựa trên model `Document.status` và UC-39
        status_mapping = {
            "DRAFT": 1, 
            "PROCESSING_WORKFLOW": 2, # (Sẽ được cập nhật bởi auto-route)
            "APPROVED": 3,
            "DELETED": 0
        }
        status_int = status_mapping.get(status_str, 1) # Mặc định là DRAFT (1)

        # 2. Ánh xạ Access Type (String) sang Access Type (Integer) và is_paid (Boolean)
        # Dựa trên model `Document.access_type` và UC-39
        access_type_mapping = {
            "public": 1,
            "paid": 1, # 'paid' vẫn là access_type=1 (public)
            "private": 4 # (Giả sử 4 là private)
            # (Bạn cần định nghĩa 2=Org, 3=Dept nếu frontend hỗ trợ)
        }
        access_type_int = access_type_mapping.get(access_type_str, 4)
        is_paid_flag = (access_type_str == 'paid') # Cờ cho UC-85
        
        document = Document(
            title=title,
            status=status_int, # Đã ánh xạ
            content=final_metadata.get('ocrContent', None), # Lưu nội dung OCR (UC-87)
            created_by=created_by,
            type=file_type,
            total_page=total_pages,
            description=description or final_metadata.get('summary', None), # Dùng summary nếu desc rỗng
            category_id=category_id,
            file_id=file_id,
            file_path=file_path,
            storage_capacity=storage_capacity,
            storage_unit=storage_unit,
            access_type=access_type_int, # Đã ánh xạ
            organization_id=user.organization_id,
            dept_id=user.dept_id,
            photo_id=photo_id,
            
            # --- Các trường mới từ Model (đã cập nhật) ---
            version=version, # "1.0"
            confidentiality=confidentiality, # "INTERNAL", "PUBLIC", "LOCKED"
            tags_json=json.dumps(final_metadata.get('tags', [])),
            is_paid=is_paid_flag # True/False
        )

        session.add(document)

        # Commit the changes to generate an id for the document
        session.commit()
        session.refresh(document)  # Ensure the document is bound to the session

        # Create a new UserDocument object
        user_document = UserDocument(
            user_id=user.id,
            document_id=document.id,
            type=1,  # or any other default value
            status=1,
            decentralized_by=None,
            updated_at=None,
            viewed_at=None,
            created_at=datetime.now()
        )

        session.add(user_document)

        session.commit()

        organization = session.query(Organization).filter(Organization.id == user.organization_id).first()
        if organization:
            organization.data_used = (organization.data_used or 0) + storage_capacity
            session.commit()

    return document

# *** BỔ SUNG HÀM NÀY ***
async def trigger_auto_route(document: Document, metadata: dict, user: User):
    """
    UC-39/UC-84: Tự động luân chuyển (Auto-Route)
    """
    try:
        # Lấy dữ liệu metadata đã được chốt
        category_id = metadata.get('category_id')
        confidentiality = metadata.get('confidentiality')
        
        # (Giả định trạng thái Status: 1=DRAFT, 2=PROCESSING_WORKFLOW)
        STATUS_DRAFT = 1
        STATUS_PROCESSING = 2
        
        triggered_rule = None
        
        # --- ĐỊNH NGHĨA QUY TẮC (RULES) ---
        # (Trong thực tế, nên đọc từ CSDL)
        
        # RULE 1: Báo cáo tài chính (category=1) -> Gửi cho Kế toán trưởng
        if str(category_id) == "1": # (Frontend gửi ID dạng số, nhưng form là string)
            triggered_rule = {
                "workflow_name": "Quy trình duyệt Báo cáo Tài chính",
                "process_key": "finance_report_process_v1",
                "candidate_group": "PHONG_KE_TOAN" # Nhóm người nhận
            }
            
        # RULE 2: Tài liệu mật (LOCKED) -> Gửi cho Ban An Ninh
        elif confidentiality == "LOCKED":
            triggered_rule = {
                "workflow_name": "Quy trình xử lý Tài liệu Mật",
                "process_key": "security_confidential_v1",
                "candidate_group": "BAN_AN_NINH"
            }
        
        # --- KÍCH HOẠT WORKFLOW NẾU CÓ RULE ---
        if triggered_rule:
            with session_scope() as session:
                # 1. Tạo bản ghi WorkflowInstance (UC-84)
                instance = WorkflowInstance(
                    document_id=document.id,
                    workflow_name=triggered_rule["workflow_name"],
                    process_key=triggered_rule["process_key"],
                    status="PENDING", 
                    triggered_by_user_id=user.id,
                    candidate_group=triggered_rule["candidate_group"]
                )
                session.add(instance)
                
                # 2. Cập nhật trạng thái tài liệu
                doc_to_update = session.query(Document).filter(Document.id == document.id).first()
                if doc_to_update:
                    doc_to_update.status = STATUS_PROCESSING # Chuyển từ DRAFT (1) -> PROCESSING (2)
                
                session.commit()

            logging.info(f"Đã kích hoạt Auto-Route (UC-84) cho DocID {document.id}. Chuyển đến {triggered_rule['candidate_group']}")
            
            return {
                "triggered": True,
                "workflow": { 
                    "name": triggered_rule["workflow_name"], 
                    "id": triggered_rule["process_key"] 
                },
                "message": f"Đã tự động gửi tài liệu đến: {triggered_rule['candidate_group']}.",
            }

        # Không có rule nào khớp
        logging.info(f"Không có Auto-Route (UC-84) nào khớp cho DocID {document.id}.")
        return {"triggered": False, "message": "Không có quy trình tự động nào được kích hoạt."}
        
    except Exception as e:
        logging.error(f"Lỗi khi trigger auto-route (UC-84): {e}")
        return {"triggered": False, "message": f"Lỗi auto-route: {e}"}

# *** BỔ SUNG CÁC HÀM NÀY (UC-85/86) ***
async def save_public_link(doc_id: int):
    """
    UC-86: Tạo liên kết chia sẻ công khai với thời hạn 72 giờ
    Sử dụng model PublicLink đã thêm.
    """
    with session_scope() as session:
        # Đảm bảo không tạo link trùng
        existing_link = session.query(PublicLink).filter(
            PublicLink.document_id == doc_id,
            PublicLink.is_active == True,
            PublicLink.expires_at > datetime.now()
        ).first()
        
        if existing_link:
            token = existing_link.token
        else:
            new_link = PublicLink(
                document_id=doc_id,
                # expires_at (đã được default 72h trong model)
            )
            session.add(new_link)
            session.commit()
            session.refresh(new_link)
            token = new_link.token
        
        # (Cần định nghĩa APP_DOMAIN trong config.py)
        app_domain = getattr(settings, "APP_DOMAIN", "http://localhost:5173")
        full_link = f"{app_domain}/public/view/{token}"
        return full_link

async def setup_paid_access(doc_id: int):
    """
    UC-85: Thiết lập tài liệu này cần thanh toán
    """
    with session_scope() as session:
        doc = session.query(Document).filter(Document.id == doc_id).first()
        if doc:
            doc.is_paid = True
            session.commit()
            logging.info(f"Đã thiết lập is_paid=True cho document ID: {doc_id}")
        else:
            logging.error(f"Không tìm thấy document ID: {doc_id} để thiết lập paid access.")

def delete_sentences_in_elasticsearch(document_id: int, user_id: int):
    es = Elasticsearch([elasticsearch_url])
    query = {"query": {"match": {"document_id": document_id}}}

    # Ensure meta_data index exists
    ensure_meta_data_index_exists()
    
    # Check if metadata exists
    meta_query = {
        "query": {
            "term": {
                "document_id": document_id
            }
        }
    }
    meta_response = es.search(index="meta_data", body=meta_query)
    if meta_response['hits']['total']['value'] > 0:
        # Metadata exists, delete it
        es.delete_by_query(index="meta_data", body=meta_query)

    # Determine the index based on the organization
    if check_orgnization_open_ai(user_id):
        index = "search_openai"
    else:
        index = "demo_cau"

    # Check if the document exists in the index
    doc_response = es.search(index=index, body=query)
    if doc_response['hits']['total']['value'] > 0:
        # Document exists, delete it
        es.delete_by_query(index=index, body=query)


def get_document_by_id(doc_id: int):
    with session_scope() as session:
        document = session.query(
            Document
        ).filter(
            Document.id == doc_id,
        ).first()
        if document is None:
            raise HTTPException(status_code=404, detail="Document not found!")
        return document


def serialize_document(document):
    with session_scope() as session:
        organization = session.query(Organization).filter(Organization.id == document.organization_id).first()
        category = session.query(Category).filter(Category.id == document.category_id).first()
        dept = session.query(Department).filter(Department.id == document.dept_id).first()

    meta_data = {}

    try:
        if document.type in ["pdf", "docx"]:
            # Ensure meta_data index exists
            ensure_meta_data_index_exists()
            
            es_query = {
                "query": {
                    "term": {
                        "document_id": document.id
                    }
                }
            }
            es_response = client.search(index="meta_data", body=es_query)
            if es_response["hits"]["total"]["value"] > 0:
                metadata = es_response["hits"]["hits"][0]["_source"]
                meta_data.update({
                    "document_number": metadata.get("document_number"),
                    "issuing_authority": metadata.get("issuing_authority"),
                    "date_of_issuance": metadata.get("date_of_issuance"),
                    "signature": metadata.get("signature"),
                    "agency_address": metadata.get("agency_address")
                })

        if document.type in ["jpg", "jpeg", "png"]:
            es_query = {
                "query": {
                    "term": {
                        "doc_id": document.id
                    }
                }
            }
            es_response = client.search(index="photo_info", body=es_query)
            if es_response["hits"]["total"]["value"] > 0:
                photo_info = es_response["hits"]["hits"][0]["_source"]
                meta_data.update({
                    "photo_description": photo_info.get("description")
                })
    except Exception as e:
        logging.error(f"Error in serialize_document: {e}")

    return {
        "id": document.id,
        "title": document.title,
        "category_id": document.category_id,
        "status": document.status,
        "created_by": document.created_by,
        "type": document.type,
        "total_page": document.total_page,
        "file_id": document.file_id,
        "storage_capacity": document.storage_capacity,
        "storage_unit": document.storage_unit,
        "created_at": document.created_at.isoformat() if document.created_at else None,
        "updated_at": document.updated_at.isoformat() if document.updated_at else None,
        "access_type": document.access_type,
        "organization_id": document.organization_id,
        "dept_id": document.dept_id,
        "organization_name": organization.name if organization else None,
        "category_name": category.name if category else None,
        "dept_name": dept.name if dept else None,
        "meta_data": meta_data
    }


async def get_documents_by_category(category_id: int, page: int, page_size: int, user_id: int):
    user = get_user_by_id(user_id)
    with session_scope() as session:
        if user.is_organization_manager:
            documents = session.query(
                Document
            ).join(
                Category,
                Document.category_id == Category.id
            ).filter(
                Document.category_id == category_id,
                Document.status == 1,
                Category.organization_id == user.organization_id,
                or_(
                    Document.access_type != 4,
                    Document.created_by == user.email,
                    Document.id.in_(
                        session.query(PrivateDoc.doc_id).filter(PrivateDoc.user_id == user_id)
                    )
                )
            ).offset((page - 1) * page_size).limit(page_size).all()

            total_documents = session.query(
                Document
            ).filter(
                Document.category_id == category_id,
                Document.status == 1,
                Document.organization_id == user.organization_id,
                or_(
                    Document.access_type != 4,
                    Document.created_by == user.email,
                    Document.id.in_(
                        session.query(PrivateDoc.doc_id).filter(PrivateDoc.user_id == user_id)
                    )
                )
            ).count()
        else:
            lst_user_document_private = session.query(
                PrivateDoc.doc_id,
            ).filter(
                PrivateDoc.user_id == user_id
            ).all()
            lst_user_document_private = [doc_id[0] for doc_id in lst_user_document_private]

            documents = session.query(
                Document
            ).filter(
                Document.category_id == category_id,
                Document.status == 1,
                or_(
                    and_(
                        Document.organization_id == user.organization_id,
                        Document.access_type.in_([1, 2, 3]),
                        or_(
                            Document.dept_id == user.dept_id,
                            Document.access_type == 1,
                            Document.access_type == 2,
                            and_(
                                Document.access_type == 3,
                                Document.dept_id == user.dept_id
                            )
                        )
                    ),
                    and_(
                        Document.access_type == 4,
                        or_(
                            Document.created_by == user.email,  # Uploaded by the user
                            Document.id.in_(lst_user_document_private)  # Shared with the user
                        )
                    )
                )
            ).offset((page - 1) * page_size).limit(page_size).all()

            total_documents = session.query(
                Document
            ).filter(
                Document.category_id == category_id,
                Document.status == 1,
                or_(
                    and_(
                        Document.organization_id == user.organization_id,
                        Document.access_type.in_([1, 2, 3]),
                        or_(
                            Document.dept_id == None,
                            Document.dept_id == user.dept_id,
                            Document.access_type == 1,
                            Document.access_type == 2,
                            and_(
                                Document.access_type == 3,
                                Document.dept_id == user.dept_id
                            )
                        )
                    ),
                    and_(
                        Document.access_type == 4,
                        or_(
                            Document.created_by == user.email,  # Uploaded by the user
                            Document.id.in_(lst_user_document_private)  # Shared with the user
                        )
                    )
                )
            ).count()

    return documents, total_documents


def get_user_activities(time: date, user_id: int, page: int, page_size: int):
    with session_scope() as session:
        user = session.query(User).filter(User.id == user_id).first()
        if not user:
            return [], 0

        user_documents = session.query(UserDocument).filter(
            UserDocument.user_id == user_id,
            or_(
                UserDocument.created_at >= time,
                UserDocument.updated_at >= time,
                UserDocument.viewed_at >= time
            )
        ).order_by(
            UserDocument.updated_at.desc(),
            UserDocument.created_at.desc(),
            UserDocument.viewed_at.desc()
        ).all()

        unique_documents = {}
        for user_document in user_documents:
            if user_document.document_id not in unique_documents:
                document = session.query(Document).filter(
                    Document.id == user_document.document_id,
                    Document.organization_id == user.organization_id,
                    Document.status == 1
                ).first()
                if document:
                    unique_documents[user_document.document_id] = document

        documents = list(unique_documents.values())
        total_documents = len(documents)
        start = (page - 1) * page_size
        end = start + page_size
        documents = documents[start:end]

    return documents, total_documents


def find_all_document(start, end, user_id, category_id):
    user = get_user_by_id(user_id)
    with session_scope() as session:
        documents = session.query(
            Document
        ).filter(
            Document.created_by == user.email,
            Document.created_at >= start,
            Document.created_at <= end,
            Document.category_id == category_id,
            Document.status == 1,
        ).all()
        return documents


def get_document_collection_size():
    with session_scope() as session:
        documents = session.query(Document).count()
        return documents


def search(query, client, category_id=None, type_doc=None, user_id=None):
    if check_orgnization_open_ai(user_id):
        index = "search_openai"
        headers = {
            "Authorization": "Bearer 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",
            "Content-Type": "application/json"
        }
        query_embeddings = get_openai_embeddings_to_search([query], headers)
        query_vector = query_embeddings[0]
    else:
        index = "demo_cau"
        query_vector = embed_str([tokenize(query)])[0]

    script_query = {
        "bool": {
            "must": {
                "script_score": {
                    "query": {
                        "match_all": {}
                    },
                    "script": {
                        "source": "cosineSimilarity(params.query_vector, 'title_vector') + 1.0",
                        "params": {"query_vector": query_vector}
                    }
                }
            },
            "filter": [
                {"exists": {"field": "title_vector"}}
            ]
        }
    }
    if category_id is not None:
        script_query["bool"]["filter"].append({"term": {"category_id": category_id}})
    if type_doc is not None:
        script_query["bool"]["filter"].append({"term": {"type": type_doc}})

    response = client.search(
        index=[index],
        body={
            "query": script_query,
            "_source": {
                "includes": ["id", "title", "doc_id", "document_id"]
            },
            "sort": [
                {"_score": {"order": "desc"}}
            ]
        },
        ignore=[400]
    )
    result = []

    log = response["hits"]["hits"]
    seen_doc_ids = set()
    for hit in log:
        doc_id = hit["_source"]['doc_id']
        if doc_id not in seen_doc_ids:
            relevance_score = hit["_score"]
            relevance_percentage = relevance_score * 10
            result.append({
                "document_id": hit["_source"]['document_id'],
                "title": hit["_source"]['title'],
                "relevance_score": relevance_score,
                "relevance_percentage": relevance_percentage
            })
            seen_doc_ids.add(doc_id)

    return result, list(seen_doc_ids)


def find_custom(id_list, doc_type, start, end, user_id):
    user = get_user_by_id(user_id)
    id_list = [str(id) for id in id_list]

    with session_scope() as session:
        if user.is_organization_manager:
            list_docs_query = session.query(Document).join(
                Category,
                Document.category_id == Category.id
            ).filter(
                Document.status == 1,
                Category.organization_id == user.organization_id,
                or_(
                    Document.access_type != 4,
                    Document.created_by == user.email,
                    Document.id.in_(
                        session.query(PrivateDoc.doc_id).filter(PrivateDoc.user_id == user_id)
                    )
                )
            )
        elif user.is_dept_manager:
            list_docs_query = session.query(Document).filter(
                Document.status == 1,
                or_(
                    and_(
                        Document.organization_id == user.organization_id,
                        Document.access_type.in_([1, 2, 3]),
                        or_(
                            Document.dept_id == user.dept_id,
                            Document.access_type == 1,
                            Document.access_type == 2,
                            and_(
                                Document.access_type == 3,
                                Document.dept_id == user.dept_id
                            )
                        )
                    ),
                    and_(
                        Document.access_type == 4,
                        Document.created_by == user.email,
                        Document.id.in_(
                            session.query(PrivateDoc.doc_id).filter(PrivateDoc.user_id == user_id)
                        )
                    )
                )
            )
        else:
            lst_user_document_private = session.query(
                PrivateDoc.doc_id
            ).filter(
                PrivateDoc.user_id == user_id
            ).all()
            lst_user_document_private = [doc_id[0] for doc_id in lst_user_document_private]

            list_docs_query = session.query(Document).filter(
                Document.status == 1,
                or_(
                    and_(
                        Document.organization_id == user.organization_id,
                        Document.access_type.in_([1, 2, 3]),
                        or_(
                            Document.dept_id == user.dept_id,
                            Document.access_type == 1,
                            Document.access_type == 2,
                            and_(
                                Document.access_type == 3,
                                Document.dept_id == user.dept_id
                            )
                        )
                    ),
                    and_(
                        Document.access_type == 4,
                        Document.id.in_(lst_user_document_private)
                    )
                )
            )

        additional_docs_query = session.query(Document).filter(
            Document.status == 1,
            Document.access_type == 1,
            Document.organization_id != user.organization_id
        )

        combined_docs_query = list_docs_query.union(additional_docs_query)
        combined_docs = combined_docs_query.subquery()

        result = []
        for id in id_list:
            output = session.query(combined_docs).filter(
                and_(
                    or_(
                        combined_docs.c.documents_title == doc_type,
                        combined_docs.c.documents_content == doc_type,
                        combined_docs.c.documents_type == doc_type,
                        combined_docs.c.documents_title == id
                    ),
                    combined_docs.c.documents_status == 1
                )
            ).offset(start).limit(end).all()
            result.extend(output)

        count_query = session.query(combined_docs).filter(
            and_(
                or_(
                    combined_docs.c.documents_title == doc_type,
                    combined_docs.c.documents_content == doc_type,
                    combined_docs.c.documents_type == doc_type,
                    combined_docs.c.documents_title == id
                ),
                combined_docs.c.documents_status == 1
            )
        )
        count = count_query.count()

        return result, count


def find_custom_metadata(id_list, start, end, user_id):
    user = get_user_by_id(user_id)

    with session_scope() as session:
        if user.is_organization_manager:
            list_docs_query = session.query(Document).join(
                Category,
                Document.category_id == Category.id
            ).filter(
                Document.status == 1,
                Category.organization_id == user.organization_id,
                or_(
                    Document.access_type != 4,
                    Document.created_by == user.email,
                    Document.id.in_(
                        session.query(PrivateDoc.doc_id).filter(PrivateDoc.user_id == user_id)
                    )
                )
            )
        elif user.is_dept_manager:
            list_docs_query = session.query(Document).filter(
                Document.status == 1,
                or_(
                    and_(
                        Document.organization_id == user.organization_id,
                        Document.access_type.in_([1, 2, 3]),
                        or_(
                            Document.dept_id == user.dept_id,
                            Document.access_type == 1,
                            Document.access_type == 2,
                            and_(
                                Document.access_type == 3,
                                Document.dept_id == user.dept_id
                            )
                        )
                    ),
                    and_(
                        Document.access_type == 4,
                        Document.created_by == user.email,
                        Document.id.in_(
                            session.query(PrivateDoc.doc_id).filter(PrivateDoc.user_id == user_id)
                        )
                    )
                )
            )
        else:
            lst_user_document_private = session.query(
                PrivateDoc.doc_id
            ).filter(
                PrivateDoc.user_id == user_id
            ).all()
            lst_user_document_private = [doc_id[0] for doc_id in lst_user_document_private]

            list_docs_query = session.query(Document).filter(
                Document.status == 1,
                or_(
                    and_(
                        Document.organization_id == user.organization_id,
                        Document.access_type.in_([1, 2, 3]),
                        or_(
                            Document.dept_id == user.dept_id,
                            Document.access_type == 1,
                            Document.access_type == 2,
                            and_(
                                Document.access_type == 3,
                                Document.dept_id == user.dept_id
                            )
                        )
                    ),
                    and_(
                        Document.access_type == 4,
                        Document.id.in_(lst_user_document_private)
                    )
                )
            )

        additional_docs_query = session.query(Document).filter(
            Document.status == 1,
            Document.access_type == 1,
            Document.organization_id != user.organization_id
        )

        combined_docs_query = list_docs_query.union(additional_docs_query)
        combined_docs = combined_docs_query.subquery()

        result = []
        for id in id_list:
            output = session.query(combined_docs).filter(
                and_(
                    combined_docs.c.documents_id == id,
                    combined_docs.c.documents_status == 1
                )
            ).offset(start).limit(end).all()
            result.extend(output)

        count_query = session.query(combined_docs).filter(
            and_(
                combined_docs.c.documents_id == id,
                combined_docs.c.documents_status == 1
            )
        )
        count = count_query.count()

        return result, count


def find_id_by_title(title):
    with session_scope() as session:
        doc = session.query(Document).filter(Document.title == title).first()
        return doc.id


# def chat_with_gpt(question):
#     openai.api_key = settings.OPENAI_API_KEY
#     openai.base_url = settings.OPENAI_BASE_URL
#
#     completion = openai.chat.completions.create(
#         model="gpt-3.5-unfiltered",
#         messages=[
#             {"role": "user", "content": question},
#         ],
#     )
#
#     if completion.choices:
#         return completion.choices[0].message.content
#     else:
#         return None

def chat_with_gpt(question):
    url = "https://api.openai.com/v1/chat/completions"
    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {settings.OPENAI_API_KEY_EMBEDDING}"
    }
    data = {
        "model": "gpt-4o",
        "messages": [
            {"role": "user", "content": question}
        ]
    }

    response = requests.post(url, headers=headers, json=data)
    response_json = response.json()

    if response_json.get("choices"):
        return response_json["choices"][0]["message"]["content"]
    else:
        return None


def get_organization_by_id(user_id: int):
    with session_scope() as session:
        user = session.query(User).filter(User.id == user_id).first()
        if user and user.organization_id:
            organization = session.query(Organization).filter(Organization.id == user.organization_id).first()
            return organization
        return None


def check_orgnization_open_ai(user_id: int):
    organization = get_organization_by_id(user_id)
    if organization and organization.is_openai:
        return True
    return False


def update_category_in_elasticsearch(document_id: int, new_category_id: int, user_id: int):
    if check_orgnization_open_ai(user_id):
        for hit in scan(client, index='search_openai', query={"query": {"match": {"document_id": document_id}}}):
            client.update(index='search_openai', id=hit['_id'], body={"doc": {"category_id": new_category_id}})
    else:
        for hit in scan(client, index='demo_cau', query={"query": {"match": {"document_id": document_id}}}):
            client.update(index='demo_cau', id=hit['_id'], body={"doc": {"category_id": new_category_id}})


def process_file(upload_file: UploadFile, user_id: int):
    validate_file(upload_file)
    file_content = upload_file.file.read()
    storage_capacity = len(file_content)
    storage_unit = "bytes"
    organization = get_organization_by_id(user_id)
    limit_data = organization.limit_data if organization and organization.limit_data else 0
    data_used = organization.data_used if organization and organization.data_used else 0
    limit_token = organization.limit_token if organization and organization.limit_token else 0
    token_used = organization.token_used if organization and organization.token_used else 0
    if data_used + storage_capacity > limit_data:
        raise Exception("Adding this file would exceed the organization's data limit")
    if token_used >= limit_token:
        raise Exception("Adding this file would exceed the organization's token limit")
    return storage_capacity, storage_unit


def delete_document(doc_id: int, user_id: int):
    document = get_document_by_id(doc_id)
    user = get_user_by_id(user_id)
    if user.is_organization_manager != 1 or user.organization_id != document.organization_id:
        raise Exception("You do not have permission to move this document to trash")
    organization = get_organization_by_id(user_id)
    with session_scope() as session:
        user_documents = session.query(UserDocument).filter(UserDocument.document_id == doc_id).all()
        for user_document in user_documents:
            session.delete(user_document)
        session.commit()
        organization.data_used = organization.data_used - document.storage_capacity
        session.delete(document)
        session.commit()
    delete_sentences_in_elasticsearch(doc_id, user_id)
    delete_file_in_drive(document.file_id)
    return document


def view_document(doc_id: int, user_id: int):
    document = get_document_by_id(doc_id)
    with session_scope() as session:
        session.query(UserDocument).filter(
            UserDocument.document_id == doc_id,
            UserDocument.user_id == user_id
        ).delete()

        new_user_document = UserDocument(
            user_id=user_id,
            document_id=doc_id,
            viewed_at=datetime.now(),
            type=1,
        )
        session.add(new_user_document)
        session.commit()

        document_dict = serialize_document(document)
    return document_dict


async def get_documents_by_category_service(category_id: int, page: int, page_size: int, user_id: int):
    documents, total_documents = await get_documents_by_category(category_id, page, page_size, user_id)
    serialized_documents = [serialize_document(doc) for doc in documents]
    total_pages = (total_documents + page_size - 1) // page_size
    return serialized_documents, page, total_pages


def get_recent_documents_service(time: date, user_id: int, page: int, page_size: int):
    documents, total_documents = get_user_activities(time, user_id, page, page_size)
    total_pages = (total_documents + page_size - 1) // page_size
    return documents, page, total_pages


def serialize_datetime(obj):
    """Convert datetime objects to strings."""
    if isinstance(obj, datetime):
        return obj.isoformat()
    return obj


def search_documents(question: str, client, category_id: Union[str], type_doc: Union[str], start: int, end: int,
                     user_id: int):
    question_gpt = f"Tôi gửi 1 yêu cầu viết lại câu \"lich su VN\" trả về cho tôi dạng hoàn chỉnh của tiếng Việt \"lịch sử Việt Nam\" bỏ qua các stop word và bỏ từ 'tìm kiếm' nếu nó nằm ở đầu câu không giải thích. Hãy viết lại giúp tôi câu \"{question}\"."

    try:
        response = chat_with_gpt(question_gpt)
        if 'error' in response.lower():
            normalized_question = question
        else:
            normalized_question = response.strip('"')
    except openai.BadRequestError:
        normalized_question = question

    try:
        elasticsearch_result, doc_id_list = search(normalized_question, client, category_id, type_doc, user_id)
    except ValueError as e:
        logging.error(f"Error in search_documents: {e}")
        raise HTTPException(status_code=400, detail=str(e))

    doc_id_dict = {i['document_id']: i['title'] for i in elasticsearch_result}

    doc_id_list.insert(0, normalized_question)
    doc_id_list = list(dict.fromkeys(doc_id_list))

    list_document, list_document_size = find_custom(doc_id_list, normalized_question, start, end, user_id)

    total_page = (list_document_size + end - 1) // end

    modified_documents = []
    for doc in list_document:
        if doc is not None:
            doc_dict = document_to_dict_search(doc)
            doc_dict['related_sentences'] = doc_id_dict.get(doc_dict['id'], [])
            doc_dict['relevance_score'] = next(
                (item['relevance_score'] for item in elasticsearch_result if item['document_id'] == doc_dict['id']),
                None)
            doc_dict['relevance_percentage'] = next((item['relevance_percentage'] for item in elasticsearch_result if
                                                     item['document_id'] == doc_dict['id']), None)
            modified_documents.append(doc_dict)

    for doc in modified_documents:
        doc['created_at'] = serialize_datetime(doc.get('created_at'))
        doc['updated_at'] = serialize_datetime(doc.get('updated_at'))

    return modified_documents, list_document_size, total_page, normalized_question


def document_to_dict(document):
    return {
        "id": document.id,
        "title": document.title,
        "category_id": document.category_id,
        "status": document.status,
        "created_by": document.created_by,
        "type": document.type,
        "total_page": document.total_page,
        "description": document.description,
        "file_id": document.file_id,
        "storage_capacity": document.storage_capacity,
        "storage_unit": document.storage_unit,
        "created_at": document.created_at.isoformat() if document.created_at else None,
        "updated_at": document.updated_at.isoformat() if document.updated_at else None,
        "access_type": document.access_type,
        "organization_id": document.organization_id,
        "dept_id": document.dept_id
    }


def document_to_dict_search(document):
    doc_dict = {
        'id': getattr(document, 'documents_id', None),
        'title': getattr(document, 'documents_title', None),
        'type': getattr(document, 'documents_type', None),
        'status': getattr(document, 'documents_status', None),
        'organization_id': getattr(document, 'documents_organization_id', None),
        'dept_id': getattr(document, 'documents_dept_id', None),
        'created_by': getattr(document, 'documents_created_by', None),
        'total_page': getattr(document, 'documents_total_page', None),
        'description': getattr(document, 'documents_description', None),
        'file_id': getattr(document, 'documents_file_id', None),
        'storage_capacity': getattr(document, 'documents_storage_capacity', None),
        'storage_unit': getattr(document, 'documents_storage_unit', None),
        'created_at': getattr(document, 'documents_created_at', None),
        'updated_at': getattr(document, 'documents_updated_at', None),
        'access_type': getattr(document, 'documents_access_type', None),
        'meta_data': {}
    }

    if hasattr(document, 'documents_category_id'):
        doc_dict['category_id'] = getattr(document, 'documents_category_id', None)

    with session_scope() as session:
        if doc_dict['organization_id']:
            organization = session.query(Organization).filter(
                Organization.id == doc_dict['organization_id']
            ).first()
            doc_dict['organization_name'] = organization.name if organization else None

        if 'category_id' in doc_dict and doc_dict['category_id']:
            category = session.query(Category).filter(
                Category.id == doc_dict['category_id']
            ).first()
            doc_dict['category_name'] = category.name if category else None

        if doc_dict['dept_id']:
            dept = session.query(Department).filter(
                Department.id == doc_dict['dept_id']
            ).first()
            doc_dict['dept_name'] = dept.name if dept else None

    try:
        if doc_dict['type'] in ["pdf", "docx"]:
            es_query = {
                "query": {
                    "term": {
                        "document_id": doc_dict['id']
                    }
                }
            }
            es_response = client.search(index="meta_data", body=es_query)
            if es_response["hits"]["total"]["value"] > 0:
                metadata = es_response["hits"]["hits"][0]["_source"]
                doc_dict['meta_data'].update({
                    "document_number": metadata.get("document_number"),
                    "issuing_authority": metadata.get("issuing_authority"),
                    "date_of_issuance": metadata.get("date_of_issuance"),
                    "signature": metadata.get("signature"),
                    "agency_address": metadata.get("agency_address")
                })

        if doc_dict['type'] in ["jpg", "jpeg", "png"]:
            es_query = {
                "query": {
                    "term": {
                        "doc_id": doc_dict['id']
                    }
                }
            }
            es_response = client.search(index="photo_info", body=es_query)
            if es_response["hits"]["total"]["value"] > 0:
                photo_info = es_response["hits"]["hits"][0]["_source"]
                doc_dict['meta_data'].update({
                    "photo_description": photo_info.get("description")
                })
    except Exception as e:
        logging.error(f"Error in document_to_dict_search: {e}")

    return doc_dict

def move_to_trash(doc_id: int):
    with session_scope() as session:
        document = session.query(Document).filter(Document.id == doc_id).first()
        document.status = 0
        user_document = session.query(UserDocument).filter(UserDocument.document_id == doc_id).first()
        if user_document:
            user_document.move_to_trash_at = datetime.now()
        session.query(
            StarredDoc
        ).filter(
            StarredDoc.doc_id == doc_id
        ).delete()
        session.commit()
        document_dict = {
            "id": document.id,
            "title": document.title,
            "category_id": document.category_id,
            "status": document.status,
            "created_by": document.created_by,
            "type": document.type,
            "total_page": document.total_page,
            "description": document.description,
            "file_id": document.file_id,
            "storage_capacity": document.storage_capacity,
            "storage_unit": document.storage_unit,
            "created_at": serialize_datetime(document.created_at),
            "updated_at": serialize_datetime(document.updated_at),
            "access_type": document.access_type,
            "organization_id": document.organization_id,
            "dept_id": document.dept_id
        }
        return document_dict


def get_recent_documents_deleted_service(time: date, user_id: int, page: int, page_size: int):
    user = get_user_by_id(user_id)
    with session_scope() as session:
        if user.is_organization_manager:
            lst_document = session.query(
                Document
            ).join(
                Category,
                Document.category_id == Category.id
            ).filter(
                Document.status == 0,
                Category.organization_id == user.organization_id,
                or_(
                    Document.access_type != 4,
                    Document.created_by == user.email,
                    Document.id.in_(
                        session.query(PrivateDoc.doc_id).filter(PrivateDoc.user_id == user_id)
                    )
                )
            ).subquery()
        elif user.is_dept_manager:
            lst_document = session.query(
                Document
            ).filter(
                Document.status == 0,
                Document.dept_id == user.dept_id,
                or_(
                    Document.access_type != 4,
                    Document.created_by == user.email,
                    Document.id.in_(
                        session.query(PrivateDoc.doc_id).filter(PrivateDoc.user_id == user_id)
                    )
                )
            ).subquery()
        else:
            lst_user_document_private = session.query(
                PrivateDoc.doc_id
            ).filter(
                PrivateDoc.user_id == user_id
            ).all()
            lst_user_document_private = [doc_id[0] for doc_id in lst_user_document_private]

            lst_document = session.query(
                Document
            ).filter(
                Document.status == 0,
                or_(
                    and_(
                        Document.organization_id == user.organization_id,
                        Document.access_type.in_([1, 2, 3]),
                        or_(
                            Document.dept_id == user.dept_id,
                            Document.access_type == 1,
                            Document.access_type == 2,
                            and_(
                                Document.access_type == 3,
                                Document.dept_id == user.dept_id
                            )
                        )
                    ),
                    and_(
                        Document.access_type == 4,
                        Document.id.in_(lst_user_document_private)
                    )
                )
            ).subquery()

        user_documents = session.query(
            UserDocument
        ).join(
            lst_document,
            UserDocument.document_id == lst_document.c.id
        ).filter(
            and_(
                UserDocument.move_to_trash_at >= time,
                UserDocument.move_to_trash_at <= datetime.now(),
            )
        ).offset((page - 1) * page_size).limit(page_size).all()

        total_documents = session.query(UserDocument).filter(
            and_(
                UserDocument.move_to_trash_at >= time, UserDocument.move_to_trash_at <= datetime.now(),
                UserDocument.user_id == user_id
            )
        ).count()
        total_pages = (total_documents + page_size - 1) // page_size
        documents = []
        for user_document in user_documents:
            document = session.query(
                Document
            ).filter(
                Document.id == user_document.document_id,
                Document.status == 0,
            ).first()
            if document:
                document_data = serialize_document(document)
                document_data.update({
                    "viewed_at": serialize_datetime(user_document.viewed_at),
                    "updated_at": serialize_datetime(user_document.updated_at),
                    "created_at": serialize_datetime(user_document.created_at)
                })
                documents.append(document_data)

        return documents, page, total_pages


def permission_to_insert(user_id: int):
    with session_scope() as session:
        user = session.query(User).filter(User.id == user_id).first()
        if user:
            if user.is_organization_manager:
                return True
            else:
                return False
        else:
            return False


def get_user_by_id(user_id: int):
    with session_scope() as session:
        user = session.query(User).filter(User.id == user_id).first()
        return user


def get_category_by_id(category_id: int):
    with session_scope() as session:
        category = session.query(Category).filter(Category.id == category_id).first()
        return category


def get_category_by_doc_id(doc_id: int):
    with session_scope() as session:
        document = session.query(Document).filter(Document.id == doc_id).first()
        category = session.query(Category).filter(Category.id == document.category_id).first()
        return category


def check_access_type(access_type: int, organization_id: int, dept_id: int):
    if organization_id is not None and dept_id is not None:
        return JSONResponse(status_code=400,
                            content={"status": "400", "message": "Cannot provide both organization_id and dept_id"})
    if access_type == 2 and organization_id is None:
        return JSONResponse(status_code=400, content={"status": "400",
                                                      "message": "organization_id must be provided when access_type is 2"})
    if access_type == 3 and dept_id is None:
        return JSONResponse(status_code=400,
                            content={"status": "400", "message": "dept_id must be provided when access_type is 3"})
    return None


def recover_files(doc_id: int, user_id: int):
    with session_scope() as session:
        user = get_user_by_id(user_id)
        document = session.query(Document).filter(Document.id == doc_id).first()
        document.status = 1
        user_document = session.query(UserDocument).filter(UserDocument.document_id == doc_id,
                                                           UserDocument.user_id == user_id).first()
        if user_document:
            user_document.move_to_trash_at = None
        # Serialize the document inside the session context
        document_dict = serialize_document(document)
        session.commit()
        return document_dict


def process_upload(doc_id: int, local_file_path: str, user_id: int):
    max_retries = 5
    retries = 0
    print("Start processing document")
    try:
        while retries < max_retries:
            try:
                with session_scope() as session:
                    document = session.query(Document).filter(Document.id == doc_id).first()
                    if not document:
                        raise HTTPException(status_code=404, detail=f"Không tìm thấy tài liệu với id {doc_id}.")

                    file_format = document.title.split(".")[-1].lower()
                    if file_format in ["pdf", "docx", "txt", "xlsx", "pptx", "odt", "rtf"]:
                        try:
                            with open(local_file_path, 'rb') as file:
                                file_content = file.read()
                        except FileNotFoundError:
                            raise FileNotFoundError(f"Không tìm thấy file: {local_file_path}")

                        images = convert_file_to_content(BytesIO(file_content), document.title)
                        text = convert_images_to_text(images)
                    else:
                        text = document.description

                    document.content = text
                    session.commit()

                    # Tạo prompt cho OpenAI
                    # Kiểm tra text không None trước khi sử dụng
                    text_content = text[:1000] if text else ""
                    
                    prompt = (
                        "Trích xuất thông tin sau từ nội dung tài liệu và trả về dưới dạng JSON:\n\n"
                        "1. document_number (ví dụ: XX/AAAA/TT-ABC, XX/AAAA/NĐ-CP, ...) là giá trị của 'số:' hoặc 'luật số:' trong đó XX đại diện cho số serial của thông tư trong năm, AAAA là năm phát hành, và TT-ABC là viết tắt của cơ quan phát hành.\n"
                        "2. issuing_authority (ví dụ: Bộ Tài chính)\n"
                        "3. date_of_issuance (ví dụ: dd/mm/yyyy)\n"
                        "4. signature\n"
                        "5. agency_address (ví dụ: Hà Nội)\n\n"
                        "Vui lòng đảm bảo document_number được định dạng là XX/AAAA/TT-ABC hoặc XX/AAAA/NĐ-CP.\n\n"
                        "Nếu bất kỳ trường nào không có trong tài liệu, hãy trả về chúng là null.\n\n"
                        "Nội dung tài liệu:\n"
                        f"{text_content}"
                    )

                    response_text = call_openai(prompt)

                    if response_text:
                        try:
                            generated_data = json.loads(response_text)
                            generated_data['document_id'] = document.id
                            generated_data['category_id'] = document.category_id

                            for field in ['document_number', 'issuing_authority', 'date_of_issuance']:
                                if field not in generated_data:
                                    generated_data[field] = ""

                            if any(generated_data[field] for field in
                                   ['document_number', 'issuing_authority', 'date_of_issuance']):
                                # Ensure meta_data index exists
                                ensure_meta_data_index_exists()
                                
                                request = {
                                    "_op_type": "index",
                                    "_index": "meta_data",
                                    **generated_data
                                }

                                try:
                                    success, _ = bulk(client, [request])
                                    if not success:
                                        print("Hoạt động chỉ mục bulk không thành công.")
                                except BulkIndexError as e:
                                    print(f"Lỗi chỉ mục bulk: {e}")

                            else:
                                print("Không có trường nào cần thiết. Bỏ qua việc chỉ mục vào Elasticsearch.")

                        except json.JSONDecodeError as e:
                            print(f"Không thể phân tích JSON: {e}")

                    # organization_id mặc định là 1 để test
                    # Kiểm tra text không None trước khi tạo embeddings
                    if text and create_embeddings_and_store(text, document.title, document.category_id, document.id,
                                                             file_format, 1):
                        print("Cập nhật nội dung tài liệu thành công")
                        session.execute(delete(FileUpload).where(FileUpload.document_id == doc_id))
                        session.commit()
                        os.remove(local_file_path)

                        content = f"tài liệu {document.title} đã được xử lý hoàn tất. Giờ bạn có thể thao tác và xử lý dữ liệu với file"
                        title = f"Tài liệu {document.title} xử lý hoàn tất"
                        headers = {
                            "x-internal-api-key": "Ic0ilkVIg6nXRROAX7ytfs4zC9yfdM6Fhbgmr1bbCCmKIFnHPSXSgT4l2W58htdfhr7HyQJKRTiyycHLLhIJBsuEcLnNbOMsfD99",
                            "Authorization": "Bearer eyJhbGciOiJSUzI1NiIsInR5cCI6IkpXVCJ9.eyJhdWQiOiJ1c2VyX2NyZWRlbnRpYWxzIiwiZW1haWwiOiJhZG1pbkBnbWFpbC5jb20iLCJleHAiOjE3MjYyNDE4NDMsImlhdCI6MTcyMzY0OTg0MywiaXNzIjoiR2VuaUZhc3QtU2VhcmNoX0dvIiwibmJmIjoxNzIzNjQ5ODQzLCJ1c2VyX2lkIjoiNCJ9.hHnq9rEU_NhlJVNojhi7aZgD48vH6L_nBhZIMPhH_giv4odU6xtoPcwA0NEDp55Bc7JJhNRi4Z-3MRJj4TjFbWTJQPK8NQhvW1V4G2wEn2LDvokN0Plw65NPtKD4ARSnIvD6onOjH8ryjU1pg2X2IVZKc8pwvmHpZvavsJC4eY2_1sUVoj04fyCVBnzLg5C6ddNdXL_87oX8-mbZWqhoZlaD5-IvxBWlFBCsX3LagkceNWNxOihGi1t4KQJyFNYJcoVxb2i9bFLzbggUk2bLvJkzAj8LcnBm8g4hTnx9KlI7Q6RLYAcwwmY9qAGb4t3NCEzMeg_URpSw5OD9fMF8tMBRXvESs4UhDBCMeQfsarFC9oLM7-h2ohOF5eEglB13pmj3PFK2rHxLv2y1cyqu9UaMVh-TRCDJMzfGTj5TX_mCCHusJecZQl6aEh1C_Ta85xFeht1KjNnECKWg30C9HjxE0j6DP9t1dNvNgfXdBOjSYmugk7_yr2iPnMSOg-lwYz18dZC709l3M1QEmFmW7UZkGjwVs_WLozZ1tmvg-X_e0f8bNJUkMkToDu9DIm9pZBoST-Fz3T7y_SXyU-Yy3yYntVaOjzpTA2wG18jlfvNtqsbEv1lbWBTulxedz2HiDE4vW_Gq-RmK2prUdL8gdaen7cRES0lrMpEpwnFwjgg",
                            "Content-Type": "application/json"
                        }
                        data = {
                            "title": title,
                            "content": content,
                            "userId": user_id,
                            "docId": doc_id
                        }
                        response = requests.post("http://localhost:5173/user-notis", headers=headers, json=data)
                        response.raise_for_status()
                    return

            except (FileNotFoundError, HTTPException) as e:
                print(f"Lỗi khi xử lý tài liệu {doc_id}: {e}")
                break

            except (ValueError, SQLAlchemyError) as e:
                print(f"Lỗi khi xử lý tài liệu {doc_id}: {e}")
                break

            except Exception as e:
                retries += 1
                print(f"Lần thử lại {retries} thất bại: {e}")
                time.sleep(5)

        if retries == max_retries:
            print("Đã đạt tới số lần thử lại tối đa, công việc thất bại.")
    except Exception as e:
        print(f"Unexpected error: {e}")


def process_upload_by_openai_test(doc_id: int, local_file_path: str, user_id: int):
    max_retries = 5
    retries = 0
    orgnization_id = get_organization_by_id(user_id).id

    while retries < max_retries:
        try:
            with session_scope() as session:
                document = session.query(Document).filter(Document.id == doc_id).first()
                if not document:
                    raise HTTPException(status_code=404, detail=f"Không tìm thấy tài liệu với id {doc_id}.")

                user = session.query(User).filter(User.id == user_id).first()
                if not user:
                    raise HTTPException(status_code=404, detail=f"Không tìm thấy người dùng với id {user_id}.")

                organization = session.query(Organization).filter(Organization.id == user.organization_id).first()
                if not organization:
                    raise HTTPException(status_code=404,
                                        detail=f"Không tìm thấy tổ chức với id {user.organization_id}.")

                file_format = document.title.split(".")[-1].lower()
                if file_format in ["pdf", "docx", "txt", "xlsx", "pptx", "odt", "rtf"]:
                    if not os.path.exists(local_file_path):
                        raise FileNotFoundError(f"Không tìm thấy file: {local_file_path}")

                    with open(local_file_path, 'rb') as file:
                        file_content = file.read()

                    images = convert_file_to_content(BytesIO(file_content), document.title)
                    text = convert_images_to_text(images)
                else:
                    text = document.description

                document.content = text
                session.commit()

                prompt = (
                    "Trích xuất thông tin sau từ nội dung tài liệu và trả về dưới dạng JSON:\n\n"
                    "1. document_number (ví dụ: XX/AAAA/TT-ABC, XX/AAAA/NĐ-CP, ...) là giá trị của 'số:' hoặc 'luật số:' trong đó XX đại diện cho số serial của thông tư trong năm, AAAA là năm phát hành, và TT-ABC là viết tắt của cơ quan phát hành.\n"
                    "2. issuing_authority (ví dụ: Bộ Tài chính)\n"
                    "3. date_of_issuance (ví dụ: dd/mm/yyyy)\n"
                    "4. signature\n"
                    "5. agency_address (ví dụ: Hà Nội)\n\n"
                    "Vui lòng đảm bảo document_number được định dạng là XX/AAAA/TT-ABC hoặc XX/AAAA/NĐ-CP.\n\n"
                    "Nếu bất kỳ trường nào không có trong tài liệu, hãy trả về chúng là null.\n\n"
                    "Nội dung tài liệu:\n"
                    f"{text[:1000]}"
                )

                response_text = call_openai(prompt)

                if response_text:
                    try:
                        generated_data = json.loads(response_text)
                        generated_data['document_id'] = document.id
                        generated_data['category_id'] = document.category_id

                        #for field in ['document_number', 'issuing_authority', 'date_of_issuance']:
                        #    if field not in generated_data:
                        #        generated_data[field] = ""

                        # if any(generated_data[field] for field in
                        #        ['document_number', 'issuing_authority', 'date_of_issuance']):
                        # Ensure meta_data index exists
                        ensure_meta_data_index_exists()
                        
                        request = {
                            "_op_type": "index",
                            "_index": "meta_data",
                            **generated_data
                        }

                        try:
                            success, _ = bulk(client, [request])
                            if not success:
                                print("Hoạt động chỉ mục bulk không thành công.")
                        except BulkIndexError as e:
                            print(f"Lỗi chỉ mục bulk: {e}")

                        # else:
                        #     print("Không có trường nào cần thiết. Bỏ qua việc chỉ mục vào Elasticsearch.")

                    except json.JSONDecodeError as e:
                        print(f"Không thể phân tích JSON: {e}")

                # Kiểm tra text không None trước khi tạo embeddings
                if text and create_embeddings_and_store_openai(text, document.title, document.category_id, document.id,
                                                                file_format, orgnization_id):
                    print("Cập nhật nội dung tài liệu thành công")
                    session.execute(delete(FileUpload).where(FileUpload.document_id == doc_id))
                    session.commit()
                    os.remove(local_file_path)

                    content = f"tài liệu {document.title} đã được xử lý hoàn tất. Giờ bạn có thể thao tác và xử lý dữ liệu với file"
                    title = f"Tài liệu {document.title} xử lý hoàn tất"
                    headers = {
                        "x-internal-api-key": "Ic0ilkVIg6nXRROAX7ytfs4zC9yfdM6Fhbgmr1bbCCmKIFnHPSXSgT4l2W58htdfhr7HyQJKRTiyycHLLhIJBsuEcLnNbOMsfD99",
                        "Authorization": "Bearer eyJhbGciOiJSUzI1NiIsInR5cCI6IkpXVCJ9.eyJhdWQiOiJ1c2VyX2NyZWRlbnRpYWxzIiwiZW1haWwiOiJhZG1pbkBnbWFpbC5jb20iLCJleHAiOjE3MjYyNDE4NDMsImlhdCI6MTcyMzY0OTg0MywiaXNzIjoiR2VuaUZhc3QtU2VhcmNoX0dvIiwibmJmIjoxNzIzNjQ5ODQzLCJ1c2VyX2lkIjoiNCJ9.hHnq9rEU_NhlJVNojhi7aZgD48vH6L_nBhZIMPhH_giv4odU6xtoPcwA0NEDp55Bc7JJhNRi4Z-3MRJj4TjFbWTJQPK8NQhvW1V4G2wEn2LDvokN0Plw65NPtKD4ARSnIvD6onOjH8ryjU1pg2X2IVZKc8pwvmHpZvavsJC4eY2_1sUVoj04fyCVBnzLg5C6ddNdXL_87oX8-mbZWqhoZlaD5-IvxBWlFBCsX3LagkceNWNxOihGi1t4KQJyFNYJcoVxb2i9bFLzbggUk2bLvJkzAj8LcnBm8g4hTnx9KlI7Q6RLYAcwwmY9qAGb4t3NCEzMeg_URpSw5OD9fMF8tMBRXvESs4UhDBCMeQfsarFC9oLM7-h2ohOF5eEglB13pmj3PFK2rHxLv2y1cyqu9UaMVh-TRCDJMzfGTj5TX_mCCHusJecZQl6aEh1C_Ta85xFeht1KjNnECKWg30C9HjxE0j6DP9t1dNvNgfXdBOjSYmugk7_yr2iPnMSOg-lwYz18dZC709l3M1QEmFmW7UZkGjwVs_WLozZ1tmvg-X_e0f8bNJUkMkToDu9DIm9pZBoST-Fz3T7y_SXyU-Yy3yYntVaOjzpTA2wG18jlfvNtqsbEv1lbWBTulxedz2HiDE4vW_Gq-RmK2prUdL8gdaen7cRES0lrMpEpwnFwjgg",
                        "Content-Type": "application/json"
                    }
                    data = {
                        "title": title,
                        "content": content,
                        "userId": user_id,
                        "docId": doc_id
                    }
                    response = requests.post("http://localhost:5173/user-notis", headers=headers, json=data)
                    response.raise_for_status()
                return

        except (FileNotFoundError, HTTPException) as e:
            print(f"Lỗi khi xử lý tài liệu {doc_id}: {e}")
            break
        except (ValueError, SQLAlchemyError) as e:
            print(f"Lỗi khi xử lý tài liệu {doc_id}: {e}")
            break
        except Exception as e:
            retries += 1
            print(f"Lần thử lại {retries} thất bại: {e}")
            time.sleep(5)

    if retries == max_retries:
        print("Đã đạt tới số lần thử lại tối đa, công việc thất bại.")

# *** THAY THẾ HÀM NÀY (Thêm organization_id) ***
def create_embeddings_and_store_openai(texts: str, filename: str, category_id: int, document_id: int, type_doc: str,
                                       orgnization_id: int):
    # Xử lý trường hợp texts là None hoặc không phải là iterable
    if texts is None:
        logging.warning(f"texts is None for document {document_id}, skipping embeddings creation")
        return False
    
    text = ""
    if isinstance(texts, str):
        text = texts
    elif isinstance(texts, (list, tuple)):
        for i in texts:
            if i:
                text += str(i)
    else:
        try:
            for i in texts:
                text += str(i)
        except TypeError:
            logging.warning(f"texts is not iterable for document {document_id}, skipping embeddings creation")
            return False

    sentences = text.replace("\n", " ")

    data = sent_tokenizer.sentences_from_text(sentences)

    # Truyền organization_id vào
    result = preprocessing_indexing_elasticsearch(data, filename, category_id, document_id, type_doc, orgnization_id)
    index_batch_open_ai(result, orgnization_id)

    # Clear memory
    del data
    del result
    gc.collect()

    return True


def index_batch_open_ai(docs, organization_id):
    batch_size = 5
    headers = {
        "Authorization": "Bearer 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",
        "Content-Type": "application/json"
    }
    logging.info("Indexing documents to OpenAI...")
    total_tokens_used = 0

    for i in range(0, len(docs), batch_size):
        batch_docs = docs[i:i + batch_size]
        titles = [doc["title"] for doc in batch_docs]
        embeddings_response = get_openai_embeddings(titles, headers)
        embeddings = embeddings_response["embeddings"]
        total_tokens = embeddings_response["total_tokens"]

        total_tokens_used += total_tokens  # Accumulate total tokens

        requests = []
        for j, doc in enumerate(batch_docs):
            if j < len(embeddings):
                request = {
                    "_op_type": "index",
                    "_index": "search_openai",
                    "title_vector": embeddings[j],
                    **doc
                }
                requests.append(request)
            else:
                logging.error(f"Embedding for document {j} not found.")
                continue

        try:
            success, _ = bulk(client, requests)
            if not success:
                logging.error("Bulk index operation did not succeed.")
        except BulkIndexError as e:
            logging.error(f"Bulk index error: {e}")

        del requests
        del titles
        del embeddings
        gc.collect()

    # Update the organization's token usage once after the loop
    with session_scope() as session:
        organization = session.query(Organization).filter(Organization.id == organization_id).first()
        if organization:
            if organization.token_used is None:
                organization.token_used = 0
            organization.token_used += total_tokens_used
            session.commit()


def get_openai_embeddings_to_search(titles, headers):
    url = "https://api.openai.com/v1/embeddings"
    model = "text-embedding-ada-002"
    encoding_format = "float"

    embeddings = []
    for title in titles:
        payload = {
            "input": title,
            "model": model,
            "encoding_format": encoding_format
        }

        response = requests.post(url, headers=headers, json=payload)
        if response.status_code == 200:
            response_json = response.json()
            if "data" in response_json and len(response_json["data"]) > 0:
                embedding = response_json["data"][0]["embedding"]
                embeddings.append(embedding)
            else:
                print(f"Không có dữ liệu nhúng trả về cho tiêu đề: {title}")
        else:
            print(f"Yêu cầu không thành công cho tiêu đề search: {title}, mã trạng thái: {response.status_code}")

    if not embeddings:
        print("Không có nhúng nào được trả về từ OpenAI.")

    return embeddings


def get_openai_embeddings(titles, headers):
    url = "https://api.openai.com/v1/embeddings"
    model = "text-embedding-ada-002"
    encoding_format = "float"

    embeddings = []
    total_tokens = 0

    for title in titles:
        payload = {
            "input": title,
            "model": model,
            "encoding_format": encoding_format
        }

        response = requests.post(url, headers=headers, json=payload)
        if response.status_code == 200:
            response_json = response.json()
            if "data" in response_json and len(response_json["data"]) > 0:
                embedding = response_json["data"][0]["embedding"]
                embeddings.append(embedding)
                total_tokens += response_json.get("usage", {}).get("total_tokens", 0)
            else:
                print(f"Không có dữ liệu nhúng trả về cho tiêu đề: {title}")
        else:
            print(f"Yêu cầu không thành công cho tiêu đề: {title}, mã trạng thái: {response.status_code}")

    if not embeddings:
        print("Không có nhúng nào được trả về từ OpenAI.")

    return {
        "embeddings": embeddings,
        "total_tokens": total_tokens
    }


async def insert_file_upload(document_id: int, user_id: int, file_path: str):
    with session_scope() as session:
        new_file_upload = FileUpload(document_id=document_id, user_id=user_id, file_path=file_path)
        session.add(new_file_upload)
        session.commit()


def suggest_metadata(document_number=None, issuing_authority=None, date_of_issuance=None, signature=None,
                     agency_address=None):
    client = Elasticsearch([elasticsearch_url])

    query = {
        "query": {
            "bool": {
                "must": [
                    {"wildcard": {"document_number.keyword": f"*{document_number}*"}} if document_number else {
                        "match_all": {}},
                    {"wildcard": {"issuing_authority.keyword": f"*{issuing_authority}*"}} if issuing_authority else {
                        "match_all": {}},
                    {"wildcard": {"date_of_issuance.keyword": f"*{date_of_issuance}*"}} if date_of_issuance else {
                        "match_all": {}},
                    {"wildcard": {"signature.keyword": f"*{signature}*"}} if signature else {"match_all": {}},
                    {"wildcard": {"agency_address.keyword": f"*{agency_address}*"}} if agency_address else {
                        "match_all": {}}
                ]
            }
        },
        "size": 10000
    }

    response = client.search(index="meta_data", body=query)
    hits = response['hits']['hits']
    metadata = [hit['_source'] for hit in hits]

    return metadata


def filter_metadata_by_user(user_id, document_number=None, issuing_authority=None, date_of_issuance=None,
                            signature=None, agency_address=None, start=0, end=10):
    metadata_results = suggest_metadata(document_number=document_number, issuing_authority=issuing_authority,
                                        date_of_issuance=date_of_issuance, signature=signature,
                                        agency_address=agency_address)
    document_ids = [metadata['document_id'] for metadata in metadata_results]
    filtered_documents, _ = find_custom_metadata(document_ids, start, end, user_id)
    filtered_document_ids = {doc.documents_id for doc in filtered_documents}
    filtered_metadata = [metadata for metadata in metadata_results if metadata['document_id'] in filtered_document_ids]
    return filtered_metadata


def get_document_by_meta(user_id, document_number=None, issuing_authority=None, date_of_issuance=None, signature=None,
                         agency_address=None, start=0, end=10):
    metadata_results = suggest_metadata(document_number=document_number, issuing_authority=issuing_authority,
                                        date_of_issuance=date_of_issuance, signature=signature,
                                        agency_address=agency_address)
    document_ids = [metadata['document_id'] for metadata in metadata_results]
    filtered_documents, _ = find_custom_metadata(document_ids, start, end, user_id)
    modified_documents = []
    for doc in filtered_documents:
        if doc is not None:
            doc_dict = document_to_dict_search(doc)
            modified_documents.append(doc_dict)

    for doc in modified_documents:
        doc['created_at'] = serialize_datetime(doc.get('created_at'))
        doc['updated_at'] = serialize_datetime(doc.get('updated_at'))

    return modified_documents


def set_password_document(doc_id: int, password: str):
    with session_scope() as session:
        document = session.query(Document).filter(Document.id == doc_id).first()
        hashed_password = bcrypt.hashpw(password.encode('utf-8'), bcrypt.gensalt())
        document.password = hashed_password.decode('utf-8')
        session.commit()
        return


def remove_password_document(doc_id: int):
    with session_scope() as session:
        document = session.query(Document).filter(Document.id == doc_id).first()
        document.password = None
        session.commit()
        return document


def check_document_password(doc_id: int, password: str):
    with session_scope() as session:
        document = session.query(Document).filter(Document.id == doc_id).first()
        if document:
            if bcrypt.checkpw(password.encode('utf-8'), document.password.encode('utf-8')):
                return True
            else:
                return False
        else:
            return False


def check_valid_document_password(doc_id: int):
    with session_scope() as session:
        document = session.query(Document).filter(Document.id == doc_id).first()
        if document:
            if document.password:
                return True
            else:
                return False
        else:
            return False


def get_documents_by_organization_service(organization_id: int):
    document_ids = get_documents_by_organization(organization_id)
    return document_ids


def get_documents_by_organization(organization_id: int):
    with session_scope() as session:
        documents = session.query(Document.id).filter(
            Document.organization_id == organization_id,
            Document.password.isnot(None)
        ).all()

        document_ids = [doc.id for doc in documents]

    return document_ids


def get_documents_by_organization_check_image(organization_id: int):
    with session_scope() as session:
        documents = session.query(Document).filter(Document.organization_id == organization_id).all()
        return documents


def generate_image_description(doc_id: int, file_content: bytes):
    print("Starting generate_image_description")
    client = vision.ImageAnnotatorClient()
    translator = Translator()
    try:
        image = vision.Image(content=file_content)
        response = client.label_detection(image=image)
        labels = response.label_annotations
        if labels:
            description = ", ".join([label.description for label in labels])
        else:
            description = "No labels detected."
        logging.info(f"Generated description: {description}")
        if response.error.message:
            raise Exception(response.error.message)
        translated_description = translator.translate(description, dest='vi').text

        logging.info(f"Translated description: {translated_description}")
        push_to_elasticsearch(doc_id, translated_description)

        return translated_description

    except vision.exceptions.PermissionDenied as e:
        logging.error(f"PermissionDenied: {e.message}")
        raise HTTPException(status_code=403,
                            detail="Cloud Vision API is not enabled or permission is denied. Please enable the API and try again.")

    except Exception as e:
        logging.error(f"Error in generate_image_description: {e}")
        raise HTTPException(status_code=500, detail="An error occurred while processing the image.")


def push_to_elasticsearch(doc_id: int, description: str):
    try:
        # Check if the index exists
        if not client.indices.exists(index="photo_info"):
            # Create the index with mappings
            client.indices.create(
                index="photo_info",
                body={
                    "mappings": {
                        "properties": {
                            "doc_id": {"type": "integer"},
                            "description": {"type": "text"}
                        }
                    }
                }
            )

        document = {
            "doc_id": doc_id,
            "description": description
        }
        response = client.index(index="photo_info", id=doc_id, body=document)
        logging.info(f"Document indexed in Elasticsearch: {response}")
    except Exception as e:
        logging.error(f"Error indexing document in Elasticsearch: {e}")
        raise HTTPException(status_code=500, detail="Failed to index document in Elasticsearch.")


def search_image_descriptions(query: str, user_id: int):
    user = get_user_by_id(user_id)
    with session_scope() as session:
        try:
            # Define the search query
            search_query = {
                "query": {
                    "match": {
                        "description": query
                    }
                }
            }

            # Execute the search query
            response = client.search(index="photo_info", body=search_query)

            # Extract the search results
            results = response['hits']['hits']
            doc_ids = [result["_source"]["doc_id"] for result in results]

            # Filter doc_ids based on user permissions
            if user.is_organization_manager:
                accessible_docs = session.query(Document).join(
                    Category,
                    Document.category_id == Category.id
                ).filter(
                    Document.id.in_(doc_ids),
                    Document.status == 1,
                    Category.organization_id == user.organization_id,
                    or_(
                        Document.access_type != 4,
                        Document.created_by == user.email,
                        Document.id.in_(
                            session.query(PrivateDoc.doc_id).filter(PrivateDoc.user_id == user_id)
                        )
                    )
                ).all()
            else:
                lst_user_document_private = session.query(
                    PrivateDoc.doc_id,
                ).filter(
                    PrivateDoc.user_id == user_id
                ).all()
                lst_user_document_private = [doc_id[0] for doc_id in lst_user_document_private]

                accessible_docs = session.query(Document).filter(
                    Document.id.in_(doc_ids),
                    Document.status == 1,
                    or_(
                        and_(
                            Document.organization_id == user.organization_id,
                            Document.access_type.in_([1, 2, 3]),
                            or_(
                                Document.dept_id == user.dept_id,
                                Document.access_type == 1,
                                Document.access_type == 2,
                                and_(
                                    Document.access_type == 3,
                                    Document.dept_id == user.dept_id
                                )
                            )
                        ),
                        and_(
                            Document.access_type == 4,
                            or_(
                                Document.created_by == user.email,  # Uploaded by the user
                                Document.id.in_(lst_user_document_private)  # Shared with the user
                            )
                        )
                    )
                ).all()

            return accessible_docs

        except Exception as e:
            logging.error(f"Error searching image descriptions: {e}")
            raise HTTPException(status_code=500, detail="An error occurred while searching for image descriptions.")


# Tính toán cosine similarity
def rabin_karp(text, pattern):
    d = 256  # Number of characters in the input alphabet
    q = 101  # A prime number
    m = len(pattern)
    n = len(text)
    p = 0
    t = 0
    h = 1
    result = []

    if m == 0 or n == 0 or m > n:
        return result

    for i in range(m - 1):
        h = (h * d) % q

    for i in range(m):
        p = (d * p + ord(pattern[i])) % q
        t = (d * t + ord(text[i])) % q

    for i in range(n - m + 1):
        if p == t:
            for j in range(m):
                if text[i + j] != pattern[j]:
                    break
            j += 1
            if j == m:
                result.append(i)

        # Calculate hash value for next window of text: Remove leading digit, add trailing digit
        if i < n - m:
            t = (d * (t - ord(text[i]) * h) + ord(text[i + m])) % q
            if t < 0:
                t = t + q

    return result


def advanced_rabin_karp_similarity(text1, text2):
    # Convert numpy.float64 to string
    text1 = ''.join(map(str, text1))
    text2 = ''.join(map(str, text2))
    matches = rabin_karp(text1, text2)
    return len(matches) / max(len(text1), len(text2))


def advanced_cosine_similarity(a, b):
    norm_a = np.linalg.norm(a)
    norm_b = np.linalg.norm(b)
    return np.dot(a, b) / (norm_a * norm_b) if norm_a > 0 and norm_b > 0 else 0


def normalize_vector(vector):
    norm = np.linalg.norm(vector)
    return vector / norm if norm > 0 else vector


def chunk_document_by_paragraph(text):
    paragraphs = text.split('\n\n')  # Giả sử các đoạn văn cách nhau bằng hai dòng trắng
    return [paragraph for paragraph in paragraphs if paragraph.strip() != '']


def chunk_document(text, max_tokens=8192):
    words = text.split()
    chunks = []
    current_chunk = []
    current_length = 0

    for word in words:
        word_length = len(word)
        if current_length + word_length + 1 > max_tokens:
            chunks.append(' '.join(current_chunk))
            current_chunk = [word]
            current_length = word_length
        else:
            current_chunk.append(word)
            current_length += word_length + 1

    if current_chunk:
        chunks.append(' '.join(current_chunk))

    return chunks


def get_openai_embeddings_check_plagiarism(texts, model="text-embedding-ada-002"):
    embeddings = []
    for text in texts:
        try:
            response = openai.Embedding.create(
                input=[text],
                model=model
            )
            embeddings.extend([normalize_vector(np.array(embedding['embedding'])) for embedding in response['data']])
        except openai.error.APIError as e:
            logger.error(f"OpenAI API error: {e}")
            return []
    return embeddings


def calculate_average_similarity(doc_embeddings, other_embeddings):
    similarities = []
    for doc_embedding in doc_embeddings:
        for other_embedding in other_embeddings:
            similarity = advanced_cosine_similarity(doc_embedding, other_embedding)
            similarities.append(similarity)
    return np.mean(similarities)


def describe_similarity(similarity):
    if similarity >= 0.90:
        return "Có khả năng đạo văn cao."
    elif similarity >= 0.85:
        return "Xem xét kỹ hơn, nhưng không chắc chắn là đạo văn."
    else:
        return "Thường ít khả năng là đạo văn."


def check_plagiarism(doc_id: int):
    valid_extensions = ['docx', 'pdf', 'xlsx', 'pptx']

    with session_scope() as session:
        document = session.query(Document).filter(Document.id == doc_id).first()
        if not document:
            raise HTTPException(status_code=404, detail="Document not found")

        if not any(document.file_path.endswith(ext) for ext in valid_extensions):
            raise HTTPException(status_code=400, detail="Invalid document type")

        organization_id = document.organization_id
        other_documents = session.query(Document).filter(
            Document.organization_id == organization_id,
            Document.id != doc_id
        ).all()

        other_documents = [doc for doc in other_documents if
                           any(doc.file_path.endswith(ext) for ext in valid_extensions)]

        doc_chunks = chunk_document(document.content)
        doc_embeddings = get_openai_embeddings_check_plagiarism(doc_chunks)

        similarity_map = {}

        for other_doc in other_documents:
            other_chunks = chunk_document(other_doc.content)
            other_embeddings = get_openai_embeddings_check_plagiarism(other_chunks)

            avg_similarity = calculate_average_similarity(doc_embeddings, other_embeddings)

            if other_doc.id not in similarity_map or avg_similarity > similarity_map[other_doc.id]["similarity"]:
                similarity_map[other_doc.id] = {
                    "similarity": avg_similarity,
                    "document": serialize_document(other_doc)
                }

        top_similarities = [
            {
                "document": details["document"],
                "similarity": details["similarity"],
                "description": describe_similarity(details["similarity"]),
            }
            for doc_id, details in similarity_map.items()
        ]

        top_similarities.sort(key=lambda x: x['similarity'], reverse=True)

        top_5_similarities = top_similarities[:5]

        return top_5_similarities


def calculate_text_similarity(text1, text2):
    from sklearn.feature_extraction.text import TfidfVectorizer
    from sklearn.metrics.pairwise import cosine_similarity
    import logging

    # Check if texts are empty or contain only stop words
    if not text1.strip() or not text2.strip():
        logging.error("One or both texts are empty or contain only stop words")
        raise ValueError("One or both texts are empty or contain only stop words")

    vectorizer = TfidfVectorizer()
    try:
        tfidf_matrix = vectorizer.fit_transform([text1, text2])
        if tfidf_matrix.shape[1] == 0:
            logging.error("Empty vocabulary; perhaps the documents only contain stop words")
            raise ValueError("Empty vocabulary; perhaps the documents only contain stop words")
        similarity_matrix = cosine_similarity(tfidf_matrix[0:1], tfidf_matrix[1:2])
        return similarity_matrix[0][0]
    except ValueError as e:
        logging.error(f"Error calculating text similarity: {e}")
        raise


def calculate_jaccard_similarity(text1, text2):
    # Tạo vector n-grams từ văn bản
    vectorizer = CountVectorizer(ngram_range=(1, 2))  # Sử dụng n-grams
    X = vectorizer.fit_transform([text1, text2])

    # Chuyển ma trận sparse thành ma trận dense để thực hiện các phép toán
    X_dense = X.toarray()

    # Tính toán sự tương đồng Jaccard
    intersection = np.sum(np.minimum(X_dense[0, :], X_dense[1, :]))
    union = np.sum(np.maximum(X_dense[0, :], X_dense[1, :]))

    return intersection / union if union > 0 else 0


def compare_document_similarity(doc_id1: int, doc_id2: int):
    valid_extensions = ['docx', 'pdf', 'xlsx', 'pptx']

    with session_scope() as session:
        document1 = session.query(Document).filter(Document.id == doc_id1).first()
        document2 = session.query(Document).filter(Document.id == doc_id2).first()

        if not document1 or not document2:
            raise HTTPException(status_code=404, detail="Một hoặc cả hai tài liệu không được tìm thấy")

        if not any(document1.title.endswith(ext) for ext in valid_extensions) or \
                not any(document2.title.endswith(ext) for ext in valid_extensions):
            raise HTTPException(status_code=400, detail="Loại tài liệu không hợp lệ")

        doc_chunks1 = chunk_document(document1.content)
        doc_chunks2 = chunk_document(document2.content)

        doc_embeddings1 = get_openai_embeddings_check_plagiarism(doc_chunks1)
        doc_embeddings2 = get_openai_embeddings_check_plagiarism(doc_chunks2)

        avg_embedding_similarity = calculate_average_similarity(doc_embeddings1, doc_embeddings2)
        avg_text_similarity = calculate_text_similarity(document1.content, document2.content)
        avg_jaccard_similarity = calculate_jaccard_similarity(document1.content, document2.content)

        # Trung bình trọng số của sự tương đồng nội dung, embedding, và Jaccard
        weighted_similarity = 0.5 * avg_text_similarity + 0.5 * avg_jaccard_similarity
        description = describe_similarity(weighted_similarity)

        return {
            "similarity": weighted_similarity,
            "description": description
        }


def get_image_from_db(image_id: int):
    with session_scope() as session:
        image = session.query(Document).filter(Document.id == image_id).first()
        if not image:
            raise HTTPException(status_code=404, detail=f"Image with ID {image_id} not found or is not an image")
        return image.photo_id


def create_image_from_base64(base64_string: str) -> np.ndarray:
    try:
        if isinstance(base64_string, str):
            base64_string = base64_string.encode('utf-8')
        content = base64.b64decode(base64_string)
        np_arr = np.frombuffer(content, np.uint8)
        image = cv2.imdecode(np_arr, cv2.IMREAD_COLOR)
        return image
    except (TypeError, ValueError) as e:
        print(f"Error decoding base64 string: {e}")
        return None


def calculate_ssim(image1: np.ndarray, image2: np.ndarray) -> float:
    gray_image1 = cv2.cvtColor(image1, cv2.COLOR_BGR2GRAY)
    gray_image2 = cv2.cvtColor(image2, cv2.COLOR_BGR2GRAY)
    score, _ = ssim(gray_image1, gray_image2, full=True)
    return score


def update_metadata(document_id, document_number=None, issuing_authority=None, date_of_issuance=None, signature=None,
                    agency_address=None):
    client = Elasticsearch([elasticsearch_url])

    # Ensure meta_data index exists
    ensure_meta_data_index_exists()

    query = {
        "query": {
            "term": {
                "document_id": document_id
            }
        }
    }
    response = client.search(index="meta_data", body=query)
    if response['hits']['total']['value'] == 0:
        raise HTTPException(status_code=404, detail="Document not found")

    metadata = response['hits']['hits'][0]['_source']

    if document_number is not None:
        metadata['document_number'] = document_number
    if issuing_authority is not None:
        metadata['issuing_authority'] = issuing_authority
    if date_of_issuance is not None:
        metadata['date_of_issuance'] = date_of_issuance
    if signature is not None:
        metadata['signature'] = signature
    if agency_address is not None:
        metadata['agency_address'] = agency_address

    client.update(index="meta_data", id=response['hits']['hits'][0]['_id'], body={"doc": metadata})

    return metadata


def update_photo_description(doc_id: int, new_description: str):
    client = Elasticsearch([elasticsearch_url])

    query = {
        "query": {
            "term": {
                "doc_id": doc_id
            }
        }
    }
    response = client.search(index="photo_info", body=query)
    if response['hits']['total']['value'] == 0:
        raise HTTPException(status_code=404, detail="Document not found")

    photo_info = response['hits']['hits'][0]['_source']

    photo_info['description'] = new_description

    client.update(index="photo_info", id=response['hits']['hits'][0]['_id'], body={"doc": photo_info})

    return photo_info